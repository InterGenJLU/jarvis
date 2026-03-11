"""
Document Generator Module

Generates PPTX presentations, DOCX documents, and PDF files
from structured outline data. Used by the file_editor skill
for voice-driven document creation.

Supports 14 slide types (bullets, stat, comparison, image+text, title/hero,
section divider, agenda, card grids, timeline, data table, full-bleed image,
reversed image+text, closing), bold text markup, speaker notes, slide numbers,
and theme selection.
"""

import os
import random
import re
import subprocess
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH

from core.logger import get_logger
from core.debug_logger import get_debug_logger


SHARE_DIR = Path(os.path.expanduser("~/jarvis/share"))

# Default fonts (fallback if theme doesn't specify)
HEADING_FONT = "Calibri"
BODY_FONT = "Calibri"

# ------------------------------------------------------------------
# Color themes — 10 themes, each with font pairing + accent patterns
# ------------------------------------------------------------------
THEMES = {
    "professional": {
        "heading": RGBColor(0x1B, 0x2A, 0x4A),     # Deep navy
        "subtitle": RGBColor(0x4A, 0x6F, 0xA5),     # Steel blue
        "body": RGBColor(0x2D, 0x37, 0x48),          # Dark slate
        "accent": RGBColor(0xE8, 0x91, 0x3A),        # Warm amber
        "secondary_accent": RGBColor(0xA0, 0xAE, 0xC0),  # Cool gray
        "heading_font": "Cambria",
        "body_font": "Calibri",
        "accent_patterns": ["underline", "top_bottom_bars", "corner_brackets"],
    },
    "modern": {
        "heading": RGBColor(0x0D, 0x3B, 0x3E),      # Dark teal
        "subtitle": RGBColor(0x2E, 0x8B, 0x8B),      # Medium teal
        "body": RGBColor(0x33, 0x33, 0x33),           # Charcoal
        "accent": RGBColor(0xFF, 0x6B, 0x35),         # Coral orange
        "secondary_accent": RGBColor(0xB2, 0xDF, 0xDB),  # Light teal
        "heading_font": "Century Gothic",
        "body_font": "Georgia",
        "accent_patterns": ["left_bar", "circles_cluster", "gradient_sidebar"],
    },
    "bold": {
        "heading": RGBColor(0x1A, 0x23, 0x3B),      # Dark navy
        "subtitle": RGBColor(0x55, 0x55, 0x66),      # Muted gray
        "body": RGBColor(0x33, 0x33, 0x33),           # Dark gray
        "accent": RGBColor(0xE8, 0x6C, 0x00),         # Warm orange
        "secondary_accent": RGBColor(0xFF, 0xB7, 0x4D),  # Light orange
        "heading_font": "Arial Black",
        "body_font": "Arial",
        "accent_patterns": ["top_band", "stepped_bars", "dual_tone_footer"],
    },
    "minimal": {
        "heading": RGBColor(0x1A, 0x1A, 0x2E),      # Near-black navy
        "subtitle": RGBColor(0x6B, 0x72, 0x80),      # Medium gray
        "body": RGBColor(0x37, 0x41, 0x51),           # Dark gray
        "accent": RGBColor(0x3B, 0x82, 0xF6),         # Bright blue
        "secondary_accent": RGBColor(0xE5, 0xE7, 0xEB),  # Light gray
        "heading_font": "Calibri",
        "body_font": "Calibri",
        "accent_patterns": ["underline", "line_cluster", "corner_brackets"],
    },
    "elegant": {
        "heading": RGBColor(0x21, 0x21, 0x21),      # Near black
        "subtitle": RGBColor(0x61, 0x61, 0x61),      # Medium gray
        "body": RGBColor(0x42, 0x42, 0x42),           # Dark gray
        "accent": RGBColor(0xC9, 0xA8, 0x4C),         # Muted gold
        "secondary_accent": RGBColor(0xE0, 0xE0, 0xE0),  # Silver gray
        "heading_font": "Book Antiqua",
        "body_font": "Verdana",
        "accent_patterns": ["corner_brackets", "diamond", "gradient_sidebar"],
    },
    "earth": {
        "heading": RGBColor(0x4E, 0x34, 0x2E),      # Chocolate brown
        "subtitle": RGBColor(0x8D, 0x6E, 0x63),      # Warm taupe
        "body": RGBColor(0x3E, 0x27, 0x23),           # Dark brown
        "accent": RGBColor(0xE6, 0xB3, 0x5A),         # Ochre gold
        "secondary_accent": RGBColor(0xCC, 0x59, 0x59),  # Terracotta
        "heading_font": "Georgia",
        "body_font": "Verdana",
        "accent_patterns": ["left_bar", "circles_cluster", "dual_tone_footer"],
    },
    "forest": {
        "heading": RGBColor(0x1B, 0x43, 0x32),      # Forest green
        "subtitle": RGBColor(0x55, 0x6B, 0x2F),      # Olive green
        "body": RGBColor(0x2D, 0x34, 0x36),           # Dark charcoal
        "accent": RGBColor(0xE9, 0xC4, 0x6A),         # Warm gold
        "secondary_accent": RGBColor(0xA7, 0xC4, 0xA0),  # Sage
        "heading_font": "Trebuchet MS",
        "body_font": "Georgia",
        "accent_patterns": ["stepped_bars", "dot_grid", "line_cluster"],
    },
    "ocean": {
        "heading": RGBColor(0x0C, 0x2D, 0x48),      # Deep ocean blue
        "subtitle": RGBColor(0x2E, 0x86, 0xAB),      # Cerulean
        "body": RGBColor(0x2C, 0x3E, 0x50),           # Dark blue-gray
        "accent": RGBColor(0xF1, 0x8F, 0x01),         # Tangerine
        "secondary_accent": RGBColor(0xA8, 0xDA, 0xDC),  # Ice blue
        "heading_font": "Tahoma",
        "body_font": "Cambria",
        "accent_patterns": ["top_band", "circles_cluster", "top_bottom_bars"],
    },
    "jarvis": {
        "heading": RGBColor(0x0A, 0x0E, 0x1A),      # JARVIS deep navy
        "subtitle": RGBColor(0x64, 0x74, 0x8B),      # Slate gray
        "body": RGBColor(0x1E, 0x29, 0x3B),           # Dark slate
        "accent": RGBColor(0x38, 0xBD, 0xF8),         # Signature cyan
        "secondary_accent": RGBColor(0x0C, 0x4A, 0x6E),  # Dark blue
        "heading_font": "Century Gothic",
        "body_font": "Calibri",
        "accent_patterns": ["gradient_sidebar", "dot_grid", "diamond"],
    },
    "banfield": {
        "heading": RGBColor(0x07, 0x3B, 0x4C),      # Dark teal-navy
        "subtitle": RGBColor(0x5A, 0x6E, 0x78),      # Gray-teal
        "body": RGBColor(0x1A, 0x1A, 0x1A),           # Near-black
        "accent": RGBColor(0xD7, 0x41, 0x00),         # Banfield orange
        "secondary_accent": RGBColor(0xF5, 0x7F, 0x04),  # Bright orange
        "heading_font": "Tahoma",
        "body_font": "Calibri",
        "accent_patterns": ["left_bar", "top_band", "dual_tone_footer"],
    },
}


def _hfont(theme):
    """Get heading font from theme with fallback."""
    return theme.get("heading_font", HEADING_FONT)


def _bfont(theme):
    """Get body font from theme with fallback."""
    return theme.get("body_font", BODY_FONT)


def _set_shape_opacity(shape, opacity_percent):
    """Set fill opacity for a shape. opacity_percent: 0-100 (100=fully opaque)."""
    alpha_val = str(int(opacity_percent * 1000))
    spPr = shape._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return
    color_elem = solidFill.find(qn('a:srgbClr'))
    if color_elem is None:
        color_elem = solidFill.find(qn('a:schemeClr'))
    if color_elem is None:
        return
    for existing in color_elem.findall(qn('a:alpha')):
        color_elem.remove(existing)
    alpha_elem = parse_xml(
        f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' val="{alpha_val}"/>')
    color_elem.append(alpha_elem)

_BOLD_RE = re.compile(r'\*\*(.+?)\*\*')


def _parse_bold_text(text: str) -> list:
    """Parse **bold** markdown into [(text, is_bold), ...] segments.

    Example: "**Key fact:** the detail" -> [("Key fact:", True), (" the detail", False)]
    Returns [(text, False)] if no ** markers found.
    """
    segments = []
    last_end = 0
    for match in _BOLD_RE.finditer(text):
        if match.start() > last_end:
            segments.append((text[last_end:match.start()], False))
        segments.append((match.group(1), True))
        last_end = match.end()
    if last_end < len(text):
        segments.append((text[last_end:], False))
    if not segments:
        segments.append((text, False))
    return segments


def _add_formatted_runs(paragraph, text, font_size, font_color,
                        accent_color=None, heading=False, theme=None):
    """Add text with bold/normal runs parsed from **markdown** to a paragraph.

    Bold segments render in accent_color (if provided) for visual emphasis.
    Set heading=True to use heading font instead of body font.
    Pass theme dict for theme-aware font selection.
    """
    if heading:
        font_name = _hfont(theme) if theme else HEADING_FONT
    else:
        font_name = _bfont(theme) if theme else BODY_FONT
    segments = _parse_bold_text(text)
    for seg_text, is_bold in segments:
        run = paragraph.add_run()
        run.text = seg_text
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = is_bold
        if is_bold and accent_color:
            run.font.color.rgb = accent_color
        else:
            run.font.color.rgb = font_color


class DocumentGenerator:
    """Generates PPTX, DOCX, and PDF documents from structured outlines."""

    # Dispatch table for new slide types (uniform signature)
    _SLIDE_DISPATCHERS = {
        "title_hero": "_add_title_hero_slide",
        "section_divider": "_add_section_divider_slide",
        "agenda": "_add_agenda_slide",
        "card_grid_3": "_add_card_grid_slide",
        "card_grid_4": "_add_card_grid_slide",
        "timeline": "_add_timeline_slide",
        "data_table": "_add_data_table_slide",
        "full_bleed_image": "_add_full_bleed_image_slide",
        "image_text_reversed": "_add_image_text_reversed_slide",
        "closing": "_add_closing_slide",
    }

    def __init__(self, config=None):
        self.logger = get_logger(__name__, config)
        SHARE_DIR.mkdir(parents=True, exist_ok=True)

    # ------------------------------------------------------------------
    # Accent pattern methods (12 patterns)
    # ------------------------------------------------------------------

    def _accent_underline(self, slide, prs, theme, y_pos=None):
        """Pattern: thin accent-colored bar under the title area."""
        bar_height = Inches(0.06)
        y = y_pos or Inches(1.55)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), y,
            Emu(int(prs.slide_width - Inches(1))), bar_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme["accent"]
        shape.line.fill.background()

    def _accent_left_bar(self, slide, prs, theme, **kw):
        """Pattern: vertical accent bar on left edge, full height."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.5), prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme["accent"]
        shape.line.fill.background()

    def _accent_top_band(self, slide, prs, theme, **kw):
        """Pattern: full-width accent strip at top, 0.5" tall."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), prs.slide_width, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme["accent"]
        shape.line.fill.background()

    def _accent_corner_brackets(self, slide, prs, theme, **kw):
        """Pattern: L-shaped brackets in top-left and bottom-right corners."""
        t = Inches(0.06)  # bracket thickness
        blen = Inches(1.5)  # bracket arm length
        # Top-left
        h1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), blen, t)
        v1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), t, blen)
        # Bottom-right
        sw, sh = int(prs.slide_width), int(prs.slide_height)
        h2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(sw - int(blen) - int(Inches(0.4))),
            Emu(sh - int(Inches(0.4)) - int(t)),
            blen, t)
        v2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(sw - int(Inches(0.4)) - int(t)),
            Emu(sh - int(blen) - int(Inches(0.4))),
            t, blen)
        for s in (h1, v1, h2, v2):
            s.fill.solid()
            s.fill.fore_color.rgb = theme["accent"]
            s.line.fill.background()
            _set_shape_opacity(s, 70)

    def _accent_stepped_bars(self, slide, prs, theme, **kw):
        """Pattern: 4 rectangles decreasing in width, staircase on left."""
        widths = [0.8, 0.6, 0.4, 0.2]
        opacities = [50, 40, 30, 20]
        for i, (w, op) in enumerate(zip(widths, opacities)):
            y = 1.8 + i * 1.2
            s = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(y), Inches(w), Inches(1.2))
            s.fill.solid()
            s.fill.fore_color.rgb = theme["accent"]
            s.line.fill.background()
            _set_shape_opacity(s, op)

    def _accent_gradient_sidebar(self, slide, prs, theme, **kw):
        """Pattern: 5 stacked thin rects on left, opacity fading top to bottom."""
        seg_h = int(prs.slide_height) // 5
        for i in range(5):
            s = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Emu(i * seg_h), Inches(0.25), Emu(seg_h))
            s.fill.solid()
            s.fill.fore_color.rgb = theme["accent"]
            s.line.fill.background()
            _set_shape_opacity(s, 60 - i * 12)

    def _accent_dual_tone_footer(self, slide, prs, theme, **kw):
        """Pattern: two side-by-side rects at bottom (60/40 split)."""
        sw = int(prs.slide_width)
        left_w = int(sw * 0.6)
        right_w = sw - left_w
        y = Inches(6.35)
        h = Inches(1.15)
        s1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), y, Emu(left_w), h)
        s1.fill.solid()
        s1.fill.fore_color.rgb = theme["accent"]
        s1.line.fill.background()
        s2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(left_w), y, Emu(right_w), h)
        s2.fill.solid()
        s2.fill.fore_color.rgb = theme.get("secondary_accent", theme["accent"])
        s2.line.fill.background()
        _set_shape_opacity(s2, 40)

    def _accent_circles_cluster(self, slide, prs, theme, **kw):
        """Pattern: 3 overlapping circles in bottom-right, cascading opacity."""
        specs = [
            (10.5, 5.0, 3.5, 15),
            (11.0, 5.5, 2.5, 25),
            (11.8, 6.0, 1.5, 40),
        ]
        for x, y, size, op in specs:
            s = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y), Inches(size), Inches(size))
            s.fill.solid()
            s.fill.fore_color.rgb = theme["accent"]
            s.line.fill.background()
            _set_shape_opacity(s, op)

    def _accent_top_bottom_bars(self, slide, prs, theme, **kw):
        """Pattern: thick bar at top + thin bar at bottom (asymmetric)."""
        s1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), prs.slide_width, Inches(0.35))
        s1.fill.solid()
        s1.fill.fore_color.rgb = theme["accent"]
        s1.line.fill.background()
        s2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.38), prs.slide_width, Inches(0.12))
        s2.fill.solid()
        s2.fill.fore_color.rgb = theme["accent"]
        s2.line.fill.background()
        _set_shape_opacity(s2, 40)

    def _accent_diamond(self, slide, prs, theme, **kw):
        """Pattern: rotated 45-degree square partially off right edge."""
        s = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(10.8), Inches(2.0), Inches(4.0), Inches(4.0))
        s.rotation = 45.0
        s.fill.solid()
        s.fill.fore_color.rgb = theme["accent"]
        s.line.fill.background()
        _set_shape_opacity(s, 18)

    def _accent_dot_grid(self, slide, prs, theme, **kw):
        """Pattern: 5x5 grid of small circles in top-right corner, fading."""
        dot_size = Inches(0.18)
        spacing = Inches(0.35)
        base_x, base_y = Inches(10.3), Inches(0.3)
        for row in range(5):
            for col in range(5):
                x = int(base_x) + col * int(spacing)
                y = int(base_y) + row * int(spacing)
                opacity = max(10, 50 - (row + col) * 8)
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Emu(x), Emu(y), dot_size, dot_size)
                dot.fill.solid()
                dot.fill.fore_color.rgb = theme["accent"]
                dot.line.fill.background()
                _set_shape_opacity(dot, opacity)

    def _accent_line_cluster(self, slide, prs, theme, **kw):
        """Pattern: 5 thin horizontal lines of varying lengths in bottom third."""
        lines = [
            (0.5, 5.8, 5.0),
            (1.0, 6.05, 7.5),
            (0.5, 6.3, 4.0),
            (1.5, 6.55, 8.0),
            (0.5, 6.8, 3.0),
        ]
        for x, y, w in lines:
            s = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(0.03))
            s.fill.solid()
            s.fill.fore_color.rgb = theme["accent"]
            s.line.fill.background()
            _set_shape_opacity(s, 25)

    # Accent pattern dispatcher
    _ACCENT_METHODS = {
        "underline": "_accent_underline",
        "left_bar": "_accent_left_bar",
        "top_band": "_accent_top_band",
        "corner_brackets": "_accent_corner_brackets",
        "stepped_bars": "_accent_stepped_bars",
        "gradient_sidebar": "_accent_gradient_sidebar",
        "dual_tone_footer": "_accent_dual_tone_footer",
        "circles_cluster": "_accent_circles_cluster",
        "top_bottom_bars": "_accent_top_bottom_bars",
        "diamond": "_accent_diamond",
        "dot_grid": "_accent_dot_grid",
        "line_cluster": "_accent_line_cluster",
    }

    def _apply_accent(self, slide, prs, theme):
        """Apply the accent pattern for this slide based on theme's pattern list."""
        patterns = theme.get("accent_patterns", ["underline"])
        idx = getattr(self, '_current_slide_idx', 0)
        pattern_name = patterns[idx % len(patterns)]
        method_name = self._ACCENT_METHODS.get(pattern_name, "_accent_underline")
        method = getattr(self, method_name)
        method(slide, prs, theme)

    def _add_accent_bar(self, slide, prs, theme, y_pos=None):
        """Legacy compat — delegates to _accent_underline."""
        self._accent_underline(slide, prs, theme, y_pos=y_pos)

    # ------------------------------------------------------------------
    # PPTX Generation
    # ------------------------------------------------------------------

    def create_presentation(self, structure: dict, filename: str = "presentation.pptx",
                            images: dict = None,
                            theme_name: str = "professional") -> Optional[Path]:
        """Generate a PPTX presentation from a structured outline.

        Args:
            structure: Dict with keys: title, subtitle, slides[]
                       Each slide has: title, bullets[], slide_type, notes, image_query
            filename: Output filename (saved to share/)
            images: Optional {slide_index: image_path} mapping for embedded images
            theme_name: Theme preset name (professional, modern, bold)

        Returns:
            Path to saved .pptx file, or None on failure
        """
        try:
            _dbg = get_debug_logger()
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            theme = THEMES.get(theme_name, THEMES["professional"])
            slides_data = structure.get("slides", [])
            total_slides = len(slides_data)
            h_font = _hfont(theme)
            b_font = _bfont(theme)

            _dbg.log_skill_event("doc_gen", "create_presentation_entry", {
                "filename": filename,
                "theme_name": theme_name,
                "heading_font": h_font,
                "body_font": b_font,
                "accent_patterns": theme.get("accent_patterns", ["underline"]),
                "total_slides": total_slides,
                "slide_types": [s.get("slide_type", "bullets") for s in slides_data],
                "image_indices": list(images.keys()) if images else [],
                "structure_keys": list(structure.keys()),
            })

            for i, slide_data in enumerate(slides_data):
                self._current_slide_idx = i
                slide_title = slide_data.get("title", f"Slide {i + 1}")
                bullets = slide_data.get("bullets", [])
                slide_type = slide_data.get("slide_type", "bullets")
                image_path = images.get(i) if images else None
                notes = slide_data.get("notes", "")

                # First slide: title/hero or legacy title
                if i == 0 and slide_type in ("title_hero", "bullets"):
                    dispatch_method = "_add_title_slide"
                    dispatch_reason = f"first_slide_override (type={slide_type} matched title_hero|bullets)"
                    _dbg.log_skill_event("doc_gen", "slide_dispatch", {
                        "slide_index": i,
                        "slide_type": slide_type,
                        "slide_title": slide_title,
                        "dispatch_method": dispatch_method,
                        "dispatch_reason": dispatch_reason,
                        "has_image": bool(image_path),
                        "bullet_count": len(bullets),
                        "data_keys": list(slide_data.keys()),
                    })
                    slide = self._add_title_slide(
                        prs, structure.get("title", slide_title),
                        structure.get("subtitle", ""), theme)
                # New slide types with uniform signature
                elif slide_type in self._SLIDE_DISPATCHERS:
                    dispatch_method = self._SLIDE_DISPATCHERS[slide_type]
                    dispatch_reason = "dispatch_table"
                    _dbg.log_skill_event("doc_gen", "slide_dispatch", {
                        "slide_index": i,
                        "slide_type": slide_type,
                        "slide_title": slide_title,
                        "dispatch_method": dispatch_method,
                        "dispatch_reason": dispatch_reason,
                        "has_image": bool(image_path),
                        "bullet_count": len(bullets),
                        "data_keys": list(slide_data.keys()),
                    })
                    method = getattr(self, dispatch_method)
                    slide = method(prs, slide_title, slide_data, theme,
                                   image_path=image_path)
                # Legacy: stat with validation
                elif slide_type == "stat" and slide_data.get("stat_value"):
                    dispatch_method = "_add_stat_slide"
                    dispatch_reason = "legacy_stat (has stat_value)"
                    _dbg.log_skill_event("doc_gen", "slide_dispatch", {
                        "slide_index": i,
                        "slide_type": slide_type,
                        "slide_title": slide_title,
                        "dispatch_method": dispatch_method,
                        "dispatch_reason": dispatch_reason,
                        "has_image": bool(image_path),
                        "stat_value": slide_data.get("stat_value"),
                        "data_keys": list(slide_data.keys()),
                    })
                    slide = self._add_stat_slide(prs, slide_title, slide_data, theme)
                # Legacy: comparison with validation
                elif slide_type == "comparison" and slide_data.get("left_heading"):
                    dispatch_method = "_add_comparison_slide"
                    dispatch_reason = "legacy_comparison (has left_heading)"
                    _dbg.log_skill_event("doc_gen", "slide_dispatch", {
                        "slide_index": i,
                        "slide_type": slide_type,
                        "slide_title": slide_title,
                        "dispatch_method": dispatch_method,
                        "dispatch_reason": dispatch_reason,
                        "has_image": bool(image_path),
                        "data_keys": list(slide_data.keys()),
                    })
                    slide = self._add_comparison_slide(prs, slide_title, slide_data, theme)
                # Image slide (any type with available image)
                elif image_path and Path(image_path).exists():
                    dispatch_method = "_add_image_slide"
                    dispatch_reason = f"image_fallback (type={slide_type} not in dispatchers, image exists)"
                    _dbg.log_skill_event("doc_gen", "slide_dispatch", {
                        "slide_index": i,
                        "slide_type": slide_type,
                        "slide_title": slide_title,
                        "dispatch_method": dispatch_method,
                        "dispatch_reason": dispatch_reason,
                        "image_path": str(image_path),
                        "bullet_count": len(bullets),
                        "data_keys": list(slide_data.keys()),
                    })
                    slide = self._add_image_slide(
                        prs, slide_title, bullets, image_path, theme)
                # Default: bullet content
                else:
                    dispatch_method = "_add_content_slide"
                    dispatch_reason = f"default_fallback (type={slide_type} unmatched, no image)"
                    _dbg.log_skill_event("doc_gen", "slide_dispatch", {
                        "slide_index": i,
                        "slide_type": slide_type,
                        "slide_title": slide_title,
                        "dispatch_method": dispatch_method,
                        "dispatch_reason": dispatch_reason,
                        "has_image": bool(image_path),
                        "bullet_count": len(bullets),
                        "data_keys": list(slide_data.keys()),
                    })
                    slide = self._add_content_slide(prs, slide_title, bullets, theme)

                # Speaker notes
                if notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes

                # Slide number (skip title slide)
                if i > 0:
                    self._add_slide_number(slide, i + 1, total_slides, prs)

            # Save
            output_path = SHARE_DIR / filename
            prs.save(str(output_path))
            _dbg.log_skill_event("doc_gen", "presentation_saved", {
                "filename": filename,
                "output_path": str(output_path),
                "total_slides": total_slides,
            })
            self.logger.info(f"[doc_gen] Created PPTX: {output_path} ({total_slides} slides)")
            return output_path

        except Exception as e:
            self.logger.error(f"[doc_gen] PPTX creation failed: {e}")
            return None

    def _add_title_slide(self, prs, title, subtitle, theme):
        """Add a title slide with accent strip and visual hierarchy."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_title_slide", {
            "method": "_add_title_slide (LEGACY)",
            "title": title,
            "subtitle": subtitle,
            "title_font_name": "NONE (missing font.name)",
            "subtitle_font_name": "NONE (missing font.name)",
            "title_font_size": 40,
            "subtitle_font_size": 20,
        })
        slide_layout = prs.slide_layouts[5]  # Title Only — full manual control
        slide = prs.slides.add_slide(slide_layout)

        # Clear the layout placeholder so it doesn't show default text
        # ("Click to add title") behind our custom textboxes
        if slide.placeholders:
            for ph in list(slide.placeholders):
                sp = ph._element
                sp.getparent().remove(sp)

        # Accent strip — bold bar across the bottom ~25%
        strip_height = Inches(2.0)
        strip_y = Emu(int(prs.slide_height - strip_height))
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), strip_y,
            prs.slide_width, strip_height)
        strip.fill.solid()
        strip.fill.fore_color.rgb = theme["accent"]
        strip.line.fill.background()

        # Title text — large, centered in upper area
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Emu(int(prs.slide_width - Inches(2))), Inches(2.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = theme["heading"]

        # Subtitle — inside the accent strip, white text
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Emu(int(strip_y + Inches(0.4))),
                Emu(int(prs.slide_width - Inches(2))), Inches(1.2))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            run_s = sp.add_run()
            run_s.text = subtitle
            run_s.font.name = _bfont(theme)
            run_s.font.size = Pt(20)
            run_s.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        return slide

    def _add_content_slide(self, prs, title, bullets, theme):
        """Add a content slide with title and bullet points (bold markup supported)."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_content_slide", {
            "method": "_add_content_slide",
            "title": title,
            "bullet_count": len(bullets),
            "layout_index": 1,
            "title_font_name": "NONE (uses placeholder default)",
            "body_font": "via _add_formatted_runs (BODY_FONT)",
        })
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.placeholders[0]
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = _hfont(theme)
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                _add_formatted_runs(p, bullet, Pt(18), theme["body"], theme["accent"], theme=theme)
                p.space_after = Pt(8)
                p.level = 0

        return slide

    def _add_image_slide(self, prs, title, bullets, image_path, theme):
        """Add a content slide with text on left and image on right."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_image_slide", {
            "method": "_add_image_slide",
            "title": title,
            "bullet_count": len(bullets),
            "image_path": str(image_path),
            "image_exists": Path(image_path).exists() if image_path else False,
            "title_font_name": "NONE (uses placeholder default)",
            "body_font": "via _add_formatted_runs (BODY_FONT)",
        })
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.placeholders[0]
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = _hfont(theme)
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        slide_width = prs.slide_width

        # Text box — left 55%
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8),
            Emu(int(slide_width * 0.52)), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for j, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            _add_formatted_runs(p, f"\u2022 {bullet}", Pt(16), theme["body"], theme["accent"], theme=theme)
            p.space_after = Pt(6)

        # Image — right 40%
        img_left = Emu(int(slide_width * 0.57))
        img_width = Emu(int(slide_width * 0.38))
        img_max_height = Inches(4.5)

        try:
            pic = slide.shapes.add_picture(
                str(image_path), img_left, Inches(1.8), width=img_width)
            if pic.height > img_max_height:
                ratio = img_max_height / pic.height
                pic.height = img_max_height
                pic.width = int(pic.width * ratio)
        except Exception as e:
            self.logger.warning(f"[doc_gen] Failed to add image to slide: {e}")

        return slide

    def _add_stat_slide(self, prs, title, slide_data, theme):
        """Add a stat callout slide — large centered number with label and context."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_stat_slide", {
            "method": "_add_stat_slide (LEGACY)",
            "title": title,
            "stat_value": slide_data.get("stat_value"),
            "stat_label": slide_data.get("stat_label"),
            "bullet_count": len(slide_data.get("bullets", [])),
            "title_font_name": "NONE (uses placeholder default)",
            "stat_font_name": "NONE (missing font.name)",
            "label_font_name": "NONE (missing font.name)",
        })
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.placeholders[0]
        title_shape.text = title
        for p in title_shape.text_frame.paragraphs:
            p.font.name = _hfont(theme)
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        slide_width = prs.slide_width
        content_width = Emu(int(slide_width - Inches(2)))

        # Big stat number — centered, accent color
        stat_value = slide_data.get("stat_value", "")
        stat_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.0), content_width, Inches(1.5))
        tf = stat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stat_value
        run.font.name = _hfont(theme)
        run.font.size = Pt(54)
        run.font.bold = True
        run.font.color.rgb = theme["accent"]

        # Label below the stat
        stat_label = slide_data.get("stat_label", "")
        if stat_label:
            label_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.6), content_width, Inches(0.8))
            tf2 = label_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = stat_label
            run2.font.name = _bfont(theme)
            run2.font.size = Pt(20)
            run2.font.color.rgb = theme["subtitle"]

        # Supporting context bullets
        bullets = slide_data.get("bullets", [])
        if bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.6),
                Emu(int(slide_width - Inches(3))), Inches(2.5))
            tf3 = bullet_box.text_frame
            tf3.word_wrap = True
            for j, bullet in enumerate(bullets):
                p3 = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
                _add_formatted_runs(p3, f"\u2022 {bullet}", Pt(16), theme["body"], theme["accent"], theme=theme)
                p3.space_after = Pt(4)

        return slide

    def _add_comparison_slide(self, prs, title, slide_data, theme):
        """Add a two-column comparison slide."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_comparison_slide", {
            "method": "_add_comparison_slide (LEGACY)",
            "title": title,
            "left_heading": slide_data.get("left_heading"),
            "right_heading": slide_data.get("right_heading"),
            "left_points": len(slide_data.get("left_points", [])),
            "right_points": len(slide_data.get("right_points", [])),
            "title_font_name": "NONE (uses placeholder default)",
            "heading_font_name": "NONE (missing font.name)",
        })
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.placeholders[0]
        title_shape.text = title
        for p in title_shape.text_frame.paragraphs:
            p.font.name = _hfont(theme)
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        slide_width = prs.slide_width
        col_width = Emu(int((slide_width - Inches(2)) * 0.47))
        top_y = Inches(1.8)
        col_height = Inches(4.8)

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), top_y, col_width, col_height)
        ltf = left_box.text_frame
        ltf.word_wrap = True

        left_heading = slide_data.get("left_heading", "Option A")
        p_lh = ltf.paragraphs[0]
        run_lh = p_lh.add_run()
        run_lh.text = left_heading
        run_lh.font.name = _hfont(theme)
        run_lh.font.size = Pt(22)
        run_lh.font.bold = True
        run_lh.font.color.rgb = theme["accent"]
        p_lh.space_after = Pt(12)

        for point in slide_data.get("left_points", []):
            p_l = ltf.add_paragraph()
            _add_formatted_runs(p_l, f"\u2022 {point}", Pt(16), theme["body"], theme["accent"], theme=theme)
            p_l.space_after = Pt(6)

        # Right column
        right_x = Emu(int(slide_width * 0.52))
        right_box = slide.shapes.add_textbox(
            right_x, top_y, col_width, col_height)
        rtf = right_box.text_frame
        rtf.word_wrap = True

        right_heading = slide_data.get("right_heading", "Option B")
        p_rh = rtf.paragraphs[0]
        run_rh = p_rh.add_run()
        run_rh.text = right_heading
        run_rh.font.name = _hfont(theme)
        run_rh.font.size = Pt(22)
        run_rh.font.bold = True
        run_rh.font.color.rgb = theme["accent"]
        p_rh.space_after = Pt(12)

        for point in slide_data.get("right_points", []):
            p_r = rtf.add_paragraph()
            _add_formatted_runs(p_r, f"\u2022 {point}", Pt(16), theme["body"], theme["accent"], theme=theme)
            p_r.space_after = Pt(6)

        return slide

    # ------------------------------------------------------------------
    # New slide types (uniform signature for dispatch table)
    # ------------------------------------------------------------------

    def _add_title_hero_slide(self, prs, title, slide_data, theme,
                              image_path=None):
        """Title/hero slide with large display text and optional hero image."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_title_hero_slide", {
            "method": "_add_title_hero_slide",
            "title": title,
            "subtitle": slide_data.get("subtitle", ""),
            "has_image": bool(image_path),
            "image_path": str(image_path) if image_path else None,
            "image_exists": Path(image_path).exists() if image_path else False,
            "title_font": HEADING_FONT,
            "title_size": 44,
            "subtitle_font": BODY_FONT,
            "subtitle_size": 22,
        })
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)

        # Clear default placeholders
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Optional hero image as background
        if image_path and Path(image_path).exists():
            try:
                slide.shapes.add_picture(
                    str(image_path), Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height)
                # Semi-transparent overlay for text readability
                overlay = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height)
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
                overlay.fill.fore_color.brightness = 0.0
                from pptx.oxml.ns import qn
                overlay.fill._fill.attrib[qn('a:opacity')] = '40000'
                overlay.line.fill.background()
            except Exception as e:
                self.logger.warning(f"[doc_gen] Hero image failed: {e}")

        # Accent strip at bottom
        strip_height = Inches(2.0)
        strip_y = Emu(int(prs.slide_height - strip_height))
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), strip_y,
            prs.slide_width, strip_height)
        strip.fill.solid()
        strip.fill.fore_color.rgb = theme["accent"]
        strip.line.fill.background()

        # Title — large centered
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Emu(int(prs.slide_width - Inches(2))), Inches(2.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(44)
        run.font.bold = True
        title_color = RGBColor(0xFF, 0xFF, 0xFF) if image_path else theme["heading"]
        run.font.color.rgb = title_color

        # Subtitle inside accent strip
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Emu(int(strip_y + Inches(0.4))),
                Emu(int(prs.slide_width - Inches(2))), Inches(1.2))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            run_s = sp.add_run()
            run_s.text = subtitle
            run_s.font.name = _bfont(theme)
            run_s.font.size = Pt(22)
            run_s.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        return slide

    def _add_section_divider_slide(self, prs, title, slide_data, theme,
                                    image_path=None):
        """Numbered section divider for narrative pacing."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_section_divider_slide", {
            "method": "_add_section_divider_slide",
            "title": title,
            "section_number": slide_data.get("section_number", ""),
            "subtitle": slide_data.get("subtitle", ""),
            "fonts": {"number": HEADING_FONT, "title": HEADING_FONT, "subtitle": BODY_FONT},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        section_num = slide_data.get("section_number", "")

        # Large section number — left side, accent color
        if section_num:
            num_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.0), Inches(2.5), Inches(3.0))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(section_num).zfill(2)
            run.font.name = _hfont(theme)
            run.font.size = Pt(72)
            run.font.bold = True
            run.font.color.rgb = theme["accent"]

        # Section title — to the right of the number
        title_left = Inches(4.0) if section_num else Inches(1.0)
        title_box = slide.shapes.add_textbox(
            title_left, Inches(2.2),
            Emu(int(prs.slide_width - title_left - Inches(1))), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = theme["heading"]

        # Optional subtitle
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                title_left, Inches(3.8),
                Emu(int(prs.slide_width - title_left - Inches(1))), Inches(1.0))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            run_s = sp.add_run()
            run_s.text = subtitle
            run_s.font.name = _bfont(theme)
            run_s.font.size = Pt(20)
            run_s.font.color.rgb = theme["subtitle"]

        # Bottom accent bar
        self._add_accent_bar(slide, prs, theme, y_pos=Inches(6.5))

        return slide

    def _add_agenda_slide(self, prs, title, slide_data, theme,
                          image_path=None):
        """Numbered topic list agenda slide."""
        _dbg = get_debug_logger()
        items = slide_data.get("agenda_items", slide_data.get("bullets", []))
        _dbg.log_skill_event("doc_gen", "render_agenda_slide", {
            "method": "_add_agenda_slide",
            "title": title,
            "item_count": len(items),
            "item_source": "agenda_items" if slide_data.get("agenda_items") else "bullets",
            "fonts": {"title": HEADING_FONT, "numbers": HEADING_FONT, "items": "via _add_formatted_runs"},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4),
            Emu(int(prs.slide_width - Inches(1))), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        # Agenda items — numbered list
        items = slide_data.get("agenda_items", slide_data.get("bullets", []))
        start_y = Inches(2.0)
        item_height = Inches(0.7)

        for idx, item in enumerate(items[:8]):  # Max 8 items
            y = start_y + idx * item_height

            # Number
            num_box = slide.shapes.add_textbox(
                Inches(1.0), y, Inches(1.0), item_height)
            ntf = num_box.text_frame
            np_ = ntf.paragraphs[0]
            np_.alignment = PP_ALIGN.RIGHT
            nr = np_.add_run()
            nr.text = f"{idx + 1:02d}"
            nr.font.name = _hfont(theme)
            nr.font.size = Pt(28)
            nr.font.bold = True
            nr.font.color.rgb = theme["accent"]

            # Topic text
            txt_box = slide.shapes.add_textbox(
                Inches(2.5), y, Inches(9.0), item_height)
            ttf = txt_box.text_frame
            ttf.word_wrap = True
            tp = ttf.paragraphs[0]
            _add_formatted_runs(tp, item, Pt(22), theme["body"], theme["accent"], theme=theme)

        return slide

    def _add_card_grid_slide(self, prs, title, slide_data, theme,
                             image_path=None):
        """Card grid (SmartArt substitute) — 3 or 4 columns."""
        _dbg = get_debug_logger()
        cards = slide_data.get("cards", [])
        is_4col = slide_data.get("slide_type") == "card_grid_4"
        _dbg.log_skill_event("doc_gen", "render_card_grid_slide", {
            "method": "_add_card_grid_slide",
            "title": title,
            "card_count": len(cards),
            "column_count": 4 if is_4col else 3,
            "slide_type": slide_data.get("slide_type"),
            "card_headings": [c.get("heading", "") for c in cards[:4]],
            "has_cards_key": "cards" in slide_data,
            "fonts": {"title": HEADING_FONT, "card_heading": HEADING_FONT, "card_body": "via _add_formatted_runs"},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4),
            Emu(int(prs.slide_width - Inches(1))), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        # Determine column count from slide_type
        cards = slide_data.get("cards", [])
        is_4col = slide_data.get("slide_type") == "card_grid_4"
        col_count = 4 if is_4col else 3
        cards = cards[:col_count]  # Clamp to column count

        if not cards:
            return slide

        # Layout calculations
        margin = Inches(0.5)
        gap = Inches(0.25) if is_4col else Inches(0.3)
        total_gap = gap * (col_count - 1)
        usable_width = int(prs.slide_width - 2 * margin - total_gap)
        card_width = usable_width // col_count
        card_top = Inches(1.9)
        card_height = Inches(4.5)

        for idx, card in enumerate(cards):
            x = int(margin + idx * (card_width + gap))

            # Card background — rounded rectangle
            bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, card_top, card_width, card_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            bg.line.width = Pt(1)

            # Card heading
            heading_box = slide.shapes.add_textbox(
                x + Inches(0.2), card_top + Inches(0.3),
                card_width - Inches(0.4), Inches(0.8))
            htf = heading_box.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.alignment = PP_ALIGN.CENTER
            hr = hp.add_run()
            hr.text = card.get("heading", "")
            hr.font.name = _hfont(theme)
            hr.font.size = Pt(20)
            hr.font.bold = True
            hr.font.color.rgb = theme["accent"]

            # Card description
            desc_box = slide.shapes.add_textbox(
                x + Inches(0.2), card_top + Inches(1.3),
                card_width - Inches(0.4), card_height - Inches(1.6))
            dtf = desc_box.text_frame
            dtf.word_wrap = True
            dp = dtf.paragraphs[0]
            _add_formatted_runs(
                dp, card.get("description", ""),
                Pt(14), theme["body"], theme["accent"], theme=theme)

        return slide

    def _add_timeline_slide(self, prs, title, slide_data, theme,
                            image_path=None):
        """Horizontal timeline with markers, labels, and connectors."""
        _dbg = get_debug_logger()
        points = slide_data.get("timeline_points", [])
        _dbg.log_skill_event("doc_gen", "render_timeline_slide", {
            "method": "_add_timeline_slide",
            "title": title,
            "point_count": len(points),
            "has_timeline_points_key": "timeline_points" in slide_data,
            "point_labels": [p.get("label", "") for p in points[:6]],
            "fonts": {"title": HEADING_FONT, "labels": HEADING_FONT, "descriptions": BODY_FONT},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4),
            Emu(int(prs.slide_width - Inches(1))), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        points = slide_data.get("timeline_points", [])[:6]  # Max 6 points
        if not points:
            return slide

        # Horizontal line
        line_y = Inches(3.8)
        line_left = Inches(1.5)
        line_width = Emu(int(prs.slide_width - Inches(3)))
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_left, line_y, line_width, Inches(0.04))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = theme["accent"]
        line_shape.line.fill.background()

        # Markers + labels at evenly spaced positions
        num_points = len(points)
        spacing = int(line_width) // max(num_points - 1, 1) if num_points > 1 else 0
        marker_size = Inches(0.3)

        for idx, point in enumerate(points):
            if num_points == 1:
                cx = int(line_left + line_width // 2)
            else:
                cx = int(line_left) + idx * spacing
            marker_x = cx - int(marker_size) // 2

            # Circular marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                marker_x, Emu(int(line_y - marker_size // 2)),
                marker_size, marker_size)
            marker.fill.solid()
            marker.fill.fore_color.rgb = theme["accent"]
            marker.line.fill.background()

            # Alternate labels above/below
            label_text = point.get("label", "")
            desc_text = point.get("description", "")
            label_width = Inches(2.0)
            label_x = cx - int(label_width) // 2

            if idx % 2 == 0:  # Above
                # Label
                lbox = slide.shapes.add_textbox(
                    label_x, Inches(2.2), label_width, Inches(0.5))
                ltf = lbox.text_frame
                ltf.word_wrap = True
                lp = ltf.paragraphs[0]
                lp.alignment = PP_ALIGN.CENTER
                lr = lp.add_run()
                lr.text = label_text
                lr.font.name = _hfont(theme)
                lr.font.size = Pt(14)
                lr.font.bold = True
                lr.font.color.rgb = theme["heading"]
                # Description
                dbox = slide.shapes.add_textbox(
                    label_x, Inches(2.7), label_width, Inches(1.2))
                dtf = dbox.text_frame
                dtf.word_wrap = True
                dp = dtf.paragraphs[0]
                dp.alignment = PP_ALIGN.CENTER
                dr = dp.add_run()
                dr.text = desc_text
                dr.font.name = _bfont(theme)
                dr.font.size = Pt(11)
                dr.font.color.rgb = theme["body"]
            else:  # Below
                lbox = slide.shapes.add_textbox(
                    label_x, Inches(4.3), label_width, Inches(0.5))
                ltf = lbox.text_frame
                ltf.word_wrap = True
                lp = ltf.paragraphs[0]
                lp.alignment = PP_ALIGN.CENTER
                lr = lp.add_run()
                lr.text = label_text
                lr.font.name = _hfont(theme)
                lr.font.size = Pt(14)
                lr.font.bold = True
                lr.font.color.rgb = theme["heading"]

                dbox = slide.shapes.add_textbox(
                    label_x, Inches(4.8), label_width, Inches(1.0))
                dtf = dbox.text_frame
                dtf.word_wrap = True
                dp = dtf.paragraphs[0]
                dp.alignment = PP_ALIGN.CENTER
                dr = dp.add_run()
                dr.text = desc_text
                dr.font.name = _bfont(theme)
                dr.font.size = Pt(11)
                dr.font.color.rgb = theme["body"]

        return slide

    def _add_data_table_slide(self, prs, title, slide_data, theme,
                              image_path=None):
        """Styled data table with header row and alternating shading."""
        _dbg = get_debug_logger()
        headers = slide_data.get("table_headers", [])
        rows = slide_data.get("table_rows", [])
        _dbg.log_skill_event("doc_gen", "render_data_table_slide", {
            "method": "_add_data_table_slide",
            "title": title,
            "has_table_headers": bool(headers),
            "header_count": len(headers),
            "row_count": len(rows),
            "headers": headers[:6],
            "fallback_to_bullets": not bool(headers),
            "fonts": {"title": HEADING_FONT, "header": HEADING_FONT, "cells": BODY_FONT},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4),
            Emu(int(prs.slide_width - Inches(1))), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = theme["heading"]

        self._apply_accent(slide, prs, theme)

        headers = slide_data.get("table_headers", [])
        rows = slide_data.get("table_rows", [])

        if not headers:
            # Fallback to bullet slide if no table data
            bullets = slide_data.get("bullets", [])
            for j, bullet in enumerate(bullets):
                bp = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                _add_formatted_runs(bp, bullet, Pt(16), theme["body"],
                                    theme["accent"], theme=theme)
            return slide

        # Clamp dimensions
        num_cols = min(len(headers), 6)
        num_rows = min(len(rows), 8)
        headers = headers[:num_cols]
        rows = [r[:num_cols] for r in rows[:num_rows]]

        # Add table
        table_top = Inches(1.9)
        table_left = Inches(0.8)
        table_width = Emu(int(prs.slide_width - Inches(1.6)))
        row_height = Inches(0.55)
        total_height = row_height * (num_rows + 1)

        table_shape = slide.shapes.add_table(
            num_rows + 1, num_cols,
            table_left, table_top, table_width, total_height)
        table = table_shape.table

        # Style header row
        for c, header_text in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(header_text)
            run.font.name = _hfont(theme)
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Header cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = theme["accent"]

        # Data rows with alternating shading
        alt_fill = RGBColor(0xF0, 0xF0, 0xF0)
        for r, row_data in enumerate(rows):
            for c in range(num_cols):
                cell = table.cell(r + 1, c)
                cell_text = str(row_data[c]) if c < len(row_data) else ""
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = cell_text
                run.font.name = _bfont(theme)
                run.font.size = Pt(13)
                run.font.color.rgb = theme["body"]
                # Alternating row fill
                if r % 2 == 1:
                    cell_fill = cell.fill
                    cell_fill.solid()
                    cell_fill.fore_color.rgb = alt_fill

        return slide

    def _add_full_bleed_image_slide(self, prs, title, slide_data, theme,
                                     image_path=None):
        """Full-bleed image with title overlay."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_full_bleed_image_slide", {
            "method": "_add_full_bleed_image_slide",
            "title": title,
            "overlay_text": slide_data.get("overlay_text", ""),
            "has_image": bool(image_path),
            "image_path": str(image_path) if image_path else None,
            "image_exists": Path(image_path).exists() if image_path else False,
            "fonts": {"title": HEADING_FONT, "overlay": BODY_FONT},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        has_image = image_path and Path(image_path).exists()

        if has_image:
            # Full-bleed image
            try:
                slide.shapes.add_picture(
                    str(image_path), Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height)
            except Exception as e:
                self.logger.warning(f"[doc_gen] Full-bleed image failed: {e}")
                has_image = False

        if has_image:
            # Dark overlay for text readability
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                prs.slide_width, prs.slide_height)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            overlay.line.fill.background()
            # Set 40% opacity via XML
            from pptx.oxml.ns import qn
            solid = overlay.fill._fill.find(qn('a:solidFill'))
            if solid is not None:
                clr = solid[0]
                alpha = clr.makeelement(qn('a:alpha'), {'val': '40000'})
                clr.append(alpha)
            title_color = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            # Fallback: accent-colored background
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = theme["accent"]
            bg.line.fill.background()
            title_color = RGBColor(0xFF, 0xFF, 0xFF)

        # Title — centered
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Emu(int(prs.slide_width - Inches(2))), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = title_color

        # Overlay text
        overlay_text = slide_data.get("overlay_text", "")
        if overlay_text:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5),
                Emu(int(prs.slide_width - Inches(2))), Inches(1.0))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            sr = sp.add_run()
            sr.text = overlay_text
            sr.font.name = _bfont(theme)
            sr.font.size = Pt(20)
            sr.font.color.rgb = title_color

        return slide

    def _add_image_text_reversed_slide(self, prs, title, slide_data, theme,
                                        image_path=None):
        """Image on left, text on right (mirror of standard image+text)."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_image_text_reversed_slide", {
            "method": "_add_image_text_reversed_slide",
            "title": title,
            "bullet_count": len(slide_data.get("bullets", [])),
            "has_image": bool(image_path),
            "image_path": str(image_path) if image_path else None,
            "image_exists": Path(image_path).exists() if image_path else False,
            "fonts": {"title": HEADING_FONT, "body": "via _add_formatted_runs"},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Title via placeholder
        title_shape = slide.placeholders[0]
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme["heading"]
            paragraph.font.name = _hfont(theme)

        self._apply_accent(slide, prs, theme)

        slide_width = prs.slide_width
        bullets = slide_data.get("bullets", [])

        # Image — left 40%
        if image_path and Path(image_path).exists():
            img_width = Emu(int(slide_width * 0.38))
            img_max_height = Inches(4.5)
            try:
                pic = slide.shapes.add_picture(
                    str(image_path), Inches(0.5), Inches(1.8),
                    width=img_width)
                if pic.height > img_max_height:
                    ratio = img_max_height / pic.height
                    pic.height = img_max_height
                    pic.width = int(pic.width * ratio)
            except Exception as e:
                self.logger.warning(f"[doc_gen] Reversed image failed: {e}")

        # Text — right 55%
        txt_left = Emu(int(slide_width * 0.43))
        txt_width = Emu(int(slide_width * 0.52))
        txBox = slide.shapes.add_textbox(
            txt_left, Inches(1.8), txt_width, Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for j, bullet in enumerate(bullets):
            bp = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            _add_formatted_runs(
                bp, f"\u2022 {bullet}", Pt(16), theme["body"], theme["accent"], theme=theme)
            bp.space_after = Pt(6)

        return slide

    def _add_closing_slide(self, prs, title, slide_data, theme,
                           image_path=None):
        """Closing/thank-you slide with optional contact info and takeaways."""
        _dbg = get_debug_logger()
        _dbg.log_skill_event("doc_gen", "render_closing_slide", {
            "method": "_add_closing_slide",
            "title": title,
            "closing_text": slide_data.get("closing_text", ""),
            "bullet_count": len(slide_data.get("bullets", [])),
            "fonts": {"title": HEADING_FONT, "closing": BODY_FONT, "bullets": "via _add_formatted_runs"},
        })
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Accent strip at bottom
        strip_height = Inches(1.5)
        strip_y = Emu(int(prs.slide_height - strip_height))
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), strip_y,
            prs.slide_width, strip_height)
        strip.fill.solid()
        strip.fill.fore_color.rgb = theme["accent"]
        strip.line.fill.background()

        # Title — large centered
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Emu(int(prs.slide_width - Inches(2))), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = _hfont(theme)
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = theme["heading"]

        # Closing text
        closing_text = slide_data.get("closing_text", "")
        if closing_text:
            close_box = slide.shapes.add_textbox(
                Inches(2), Inches(3.5),
                Emu(int(prs.slide_width - Inches(4))), Inches(1.0))
            ctf = close_box.text_frame
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = closing_text
            cr.font.name = _bfont(theme)
            cr.font.size = Pt(18)
            cr.font.color.rgb = theme["subtitle"]

        # Key takeaway bullets
        bullets = slide_data.get("bullets", [])
        if bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.3),
                Emu(int(prs.slide_width - Inches(4))), Inches(1.5))
            btf = bullet_box.text_frame
            btf.word_wrap = True
            for j, bullet in enumerate(bullets):
                bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                bp.alignment = PP_ALIGN.CENTER
                _add_formatted_runs(
                    bp, bullet, Pt(14), theme["body"], theme["accent"], theme=theme)
                bp.space_after = Pt(4)

        return slide

    def _add_slide_number(self, slide, number, total, prs):
        """Add slide number in the bottom-right corner."""
        num_box = slide.shapes.add_textbox(
            Emu(int(prs.slide_width - Inches(1))),
            Emu(int(prs.slide_height - Inches(0.5))),
            Inches(0.8), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{number} / {total}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ------------------------------------------------------------------
    # DOCX Generation
    # ------------------------------------------------------------------

    def create_document(self, structure: dict, filename: str = "document.docx",
                        images: dict = None) -> Optional[Path]:
        """Generate a DOCX document from a structured outline.

        Args:
            structure: Dict with keys: title, subtitle, slides[] (sections)
            filename: Output filename (saved to share/)
            images: Optional {section_index: image_path} mapping

        Returns:
            Path to saved .docx file, or None on failure
        """
        try:
            doc = Document()

            # Document title
            title_para = doc.add_heading(structure.get("title", "Document"), level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Subtitle
            subtitle = structure.get("subtitle", "")
            if subtitle:
                sub_para = doc.add_paragraph()
                sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = sub_para.add_run(subtitle)
                run.italic = True
                run.font.size = DocxPt(14)
                run.font.color.rgb = DocxRGB(0x55, 0x55, 0x77)

            doc.add_paragraph()  # Spacer

            sections = structure.get("slides", [])

            for i, section in enumerate(sections):
                if i == 0 and not section.get("bullets"):
                    continue

                section_title = section.get("title", f"Section {i}")
                bullets = section.get("bullets", [])
                slide_type = section.get("slide_type", "bullets")

                # Section heading
                doc.add_heading(section_title, level=1)

                # Stat section — prominent value + label
                if slide_type == "stat":
                    stat_val = section.get("stat_value", "")
                    stat_label = section.get("stat_label", "")
                    if stat_val:
                        stat_para = doc.add_paragraph()
                        stat_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = stat_para.add_run(stat_val)
                        run.bold = True
                        run.font.size = DocxPt(24)
                        run.font.color.rgb = DocxRGB(0x2B, 0x57, 0x9A)
                    if stat_label:
                        label_para = doc.add_paragraph()
                        label_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = label_para.add_run(stat_label)
                        run.font.size = DocxPt(12)
                        run.font.color.rgb = DocxRGB(0x66, 0x66, 0x66)

                # Image (if available)
                image_path = images.get(i) if images else None
                if image_path and Path(image_path).exists():
                    try:
                        doc.add_picture(str(image_path), width=DocxInches(5.5))
                        last_paragraph = doc.paragraphs[-1]
                        last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    except Exception as e:
                        self.logger.warning(f"[doc_gen] Failed to add image to doc section {i}: {e}")

                # Bullet points with bold markup (accent-colored leads)
                for bullet in bullets:
                    p = doc.add_paragraph(style="List Bullet")
                    segments = _parse_bold_text(bullet)
                    for seg_text, is_bold in segments:
                        run = p.add_run(seg_text)
                        run.bold = is_bold
                        if is_bold:
                            run.font.color.rgb = DocxRGB(0x2B, 0x57, 0x9A)

                # Speaker notes as callout
                notes = section.get("notes", "")
                if notes:
                    notes_para = doc.add_paragraph()
                    run = notes_para.add_run(f"\u25B6 {notes}")
                    run.italic = True
                    run.font.size = DocxPt(10)
                    run.font.color.rgb = DocxRGB(0x66, 0x66, 0x99)

            # Save
            output_path = SHARE_DIR / filename
            doc.save(str(output_path))
            self.logger.info(f"[doc_gen] Created DOCX: {output_path} ({len(sections)} sections)")
            return output_path

        except Exception as e:
            self.logger.error(f"[doc_gen] DOCX creation failed: {e}")
            return None

    # ------------------------------------------------------------------
    # PDF Conversion
    # ------------------------------------------------------------------

    def convert_to_pdf(self, source_path: Path) -> Optional[Path]:
        """Convert a PPTX or DOCX file to PDF via LibreOffice CLI.

        Tries native `libreoffice` command first, then flatpak as fallback.

        Args:
            source_path: Path to the .pptx or .docx file

        Returns:
            Path to the PDF file, or None on failure
        """
        commands = [
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(SHARE_DIR), str(source_path)],
            ["flatpak", "run", "org.libreoffice.LibreOffice",
             "--headless", "--convert-to", "pdf",
             "--outdir", str(SHARE_DIR), str(source_path)],
        ]

        for cmd in commands:
            try:
                result = subprocess.run(
                    cmd, capture_output=True, text=True, timeout=60,
                )
                if result.returncode == 0:
                    pdf_name = source_path.stem + ".pdf"
                    pdf_path = SHARE_DIR / pdf_name
                    if pdf_path.exists():
                        self.logger.info(f"[doc_gen] Converted to PDF: {pdf_path}")
                        return pdf_path
            except FileNotFoundError:
                continue
            except subprocess.TimeoutExpired:
                self.logger.error("[doc_gen] LibreOffice conversion timed out (60s)")
                return None
            except Exception as e:
                self.logger.error(f"[doc_gen] PDF conversion failed: {e}")
                continue

        self.logger.error("[doc_gen] LibreOffice not available — cannot convert to PDF")
        return None
