#!/usr/bin/env python3
"""
Presentation Engine Test Script — Full Pipeline Diagnostics

Captures THREE data streams simultaneously:
  1. Client-side: every WebSocket message sent/received with timestamps
  2. Server-side: journald logs from jarvis-web.service
  3. Debug JSONL: structured pipeline events (routing, LLM calls, skill events)
     activated via sentinel file /tmp/.jarvis_debug_active

Runs 3 test turns:
  Turn 1: Create 5-slide presentation (two-phase temp gen + new slide types)
  Turn 2: Edit slide content (pipeline cache + edit handler)
  Turn 3: Add a slide (V36 append fix — should modify, not create new file)

Post-run: inspects generated PPTX files for formatting issues (overflow, overlap).

Usage:
    python3 scripts/test_pres_engine.py [--output-dir DIR]

Outputs (all in output dir):
    pres_engine_NNN_ws_transcript.jsonl   — raw WebSocket messages
    pres_engine_NNN_debug.jsonl           — structured debug pipeline log
    pres_engine_NNN_journald.txt          — journald log capture
    pres_engine_NNN_pptx_inspect.txt      — PPTX formatting inspection
    pres_engine_NNN_summary.txt           — human-readable consolidated report
"""

import argparse
import asyncio
import json
import os
import re
import signal
import subprocess
import sys
import time
from datetime import datetime, timezone
from pathlib import Path

# ---------------------------------------------------------------------------
# PPTX inspection (embedded — no external dependency beyond python-pptx)
# ---------------------------------------------------------------------------

EMU_PER_INCH = 914400


def _emu_to_in(emu):
    if emu is None:
        return 0.0
    return round(emu / EMU_PER_INCH, 3)


def _estimate_text_overflow(text_frame, shape):
    """Return (needed_inches, shape_inches, overflows_bool)."""
    total_lines = 0
    max_fs = 12  # default pt assumption
    for para in text_frame.paragraphs:
        text = para.text
        if not text.strip():
            total_lines += 0.5
            continue
        fs = None
        for run in para.runs:
            if run.font.size:
                fs = run.font.size / 12700
                break
        if fs is None:
            fs = max_fs
        else:
            max_fs = max(max_fs, fs)
        shape_w_pt = (shape.width / EMU_PER_INCH) * 72
        chars_per_line = max(1, shape_w_pt / (0.5 * fs))
        total_lines += max(1, len(text) / chars_per_line)
    needed = (total_lines * max_fs * 1.4) / 72
    actual = shape.height / EMU_PER_INCH
    return round(needed, 2), round(actual, 2), needed > actual * 1.1


def _check_overlaps(shapes_info):
    """Return list of (shape_a, shape_b, overlap_sq_in)."""
    overlaps = []
    for i, s1 in enumerate(shapes_info):
        for j, s2 in enumerate(shapes_info):
            if j <= i:
                continue
            l1, t1, r1, b1 = s1
            l2, t2, r2, b2 = s2
            if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                area = ((min(r1, r2) - max(l1, l2)) *
                        (min(b1, b2) - max(t1, t2))) / (EMU_PER_INCH ** 2)
                overlaps.append((i, j, round(area, 3)))
    return overlaps


def inspect_pptx(filepath):
    """Inspect a PPTX file and return structured results + human text."""
    from pptx import Presentation as PptxPres

    lines = []
    issues = []

    if not os.path.exists(filepath):
        return {"error": f"File not found: {filepath}"}, f"FILE NOT FOUND: {filepath}\n"

    prs = PptxPres(filepath)
    sw, sh = prs.slide_width, prs.slide_height
    lines.append(f"FILE: {os.path.basename(filepath)}")
    lines.append(f"  Size on disk: {os.path.getsize(filepath):,} bytes")
    lines.append(f"  Slide dimensions: {_emu_to_in(sw)}\" x {_emu_to_in(sh)}\"")
    lines.append(f"  Slide count: {len(prs.slides)}")

    all_fonts = set()
    all_sizes = set()
    total_overflows = 0
    total_overlaps = 0
    slide_details = []

    for si, slide in enumerate(prs.slides):
        snum = si + 1
        lines.append(f"\n  --- Slide {snum} ---")
        try:
            lines.append(f"  Layout: {slide.slide_layout.name}")
        except Exception:
            lines.append("  Layout: unknown")

        shapes_bounds = []
        shape_names = []
        slide_overflows = 0
        slide_overlap_count = 0

        for shape in slide.shapes:
            left = shape.left or 0
            top = shape.top or 0
            width = shape.width or 0
            height = shape.height or 0
            shapes_bounds.append((left, top, left + width, top + height))
            shape_names.append(shape.name or "unnamed")

            # Out of bounds check
            oob = []
            if left < 0:
                oob.append("left<0")
            if top < 0:
                oob.append("top<0")
            if left + width > sw:
                oob.append(f"right={_emu_to_in(left+width)}\">{_emu_to_in(sw)}\"")
            if top + height > sh:
                oob.append(f"bottom={_emu_to_in(top+height)}\">{_emu_to_in(sh)}\"")

            shape_desc = (f"    [{shape.name}] type={shape.shape_type} "
                          f"pos=({_emu_to_in(left)}\",{_emu_to_in(top)}\") "
                          f"size=({_emu_to_in(width)}\"x{_emu_to_in(height)}\")")
            if oob:
                shape_desc += f" *** OOB: {', '.join(oob)}"
                issues.append(f"Slide {snum}: {shape.name} out of bounds — {', '.join(oob)}")

            # Text overflow check
            if shape.has_text_frame:
                tf = shape.text_frame
                needed, actual, overflows = _estimate_text_overflow(tf, shape)
                text_preview = tf.text[:80].replace('\n', ' ')
                if overflows and actual > 0.1:  # skip decorative thin shapes
                    slide_overflows += 1
                    total_overflows += 1
                    shape_desc += f" *** OVERFLOW: needs ~{needed}\" has {actual}\""
                    issues.append(
                        f"Slide {snum}: {shape.name} text overflow — "
                        f"needs ~{needed}\" in {actual}\" container "
                        f"(\"{text_preview[:60]}...\")")

                # Collect font info
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            all_fonts.add(run.font.name)
                        if run.font.size:
                            all_sizes.add(round(run.font.size / 12700, 1))

            lines.append(shape_desc)

        # Check overlaps on this slide
        overlaps = _check_overlaps(shapes_bounds)
        if overlaps:
            for a, b, area in overlaps:
                if area > 0.5:  # only report significant overlaps
                    total_overlaps += 1
                    slide_overlap_count += 1
                    issues.append(
                        f"Slide {snum}: {shape_names[a]} <-> {shape_names[b]} "
                        f"overlap {area} sq in")
            if slide_overlap_count:
                lines.append(f"    *** {slide_overlap_count} significant overlap(s) on this slide")

        slide_details.append({
            "slide_num": snum,
            "shape_count": len(shapes_bounds),
            "overflows": slide_overflows,
            "overlaps": slide_overlap_count,
        })

    # Summary
    lines.append(f"\n  SUMMARY: {len(prs.slides)} slides, "
                 f"{sum(s['shape_count'] for s in slide_details)} shapes")
    lines.append(f"  Fonts: {sorted(all_fonts) if all_fonts else 'inherited'}")
    lines.append(f"  Font sizes: {sorted(all_sizes) if all_sizes else 'inherited'}")
    lines.append(f"  Text overflows: {total_overflows}")
    lines.append(f"  Significant overlaps: {total_overlaps}")
    if issues:
        lines.append(f"\n  ALL ISSUES ({len(issues)}):")
        for iss in issues:
            lines.append(f"    - {iss}")

    result = {
        "file": os.path.basename(filepath),
        "slide_count": len(prs.slides),
        "total_shapes": sum(s['shape_count'] for s in slide_details),
        "fonts": sorted(all_fonts),
        "font_sizes": sorted(all_sizes),
        "overflows": total_overflows,
        "overlaps": total_overlaps,
        "issues": issues,
        "slides": slide_details,
    }
    return result, "\n".join(lines)


# ---------------------------------------------------------------------------
# WebSocket client
# ---------------------------------------------------------------------------

async def ws_send_and_collect(uri, message_text, turn_label, transcript,
                              timeout=300):
    """Send a chat message over WS and collect all response messages.

    Returns dict with response text, timing, and all raw WS messages.
    """
    import aiohttp

    turn_record = {
        "turn": turn_label,
        "sent": message_text,
        "sent_ts": time.time(),
        "sent_iso": datetime.now(timezone.utc).isoformat(),
        "ws_messages_received": [],
        "response_text": None,
        "stats": None,
        "error": None,
        "elapsed_s": None,
    }

    try:
        async with aiohttp.ClientSession() as session:
            async with session.ws_connect(uri, timeout=30) as ws:
                # Drain the initial burst (history, session_list, system_stats)
                initial_messages = []
                try:
                    while True:
                        msg = await asyncio.wait_for(ws.receive(), timeout=3)
                        if msg.type == aiohttp.WSMsgType.TEXT:
                            data = json.loads(msg.data)
                            initial_messages.append({
                                "ts": time.time(),
                                "direction": "recv_initial",
                                "type": data.get("type"),
                                "size": len(msg.data),
                            })
                            # system_stats is typically the last initial message
                            if data.get("type") == "system_stats":
                                break
                        elif msg.type in (aiohttp.WSMsgType.CLOSED,
                                          aiohttp.WSMsgType.ERROR):
                            break
                except asyncio.TimeoutError:
                    pass  # done draining

                turn_record["initial_messages"] = initial_messages

                # Send our message
                payload = json.dumps({"type": "message", "content": message_text})
                await ws.send_str(payload)
                transcript.append({
                    "ts": time.time(),
                    "direction": "sent",
                    "turn": turn_label,
                    "payload": {"type": "message", "content": message_text},
                })

                t0 = time.time()

                # Collect responses until we see the 'stats' message
                # (which is the end-of-response signal per _handle_chat_message)
                got_response = False
                got_stats = False
                while True:
                    remaining = timeout - (time.time() - t0)
                    if remaining <= 0:
                        turn_record["error"] = f"Timeout after {timeout}s"
                        break
                    try:
                        msg = await asyncio.wait_for(
                            ws.receive(),
                            timeout=min(remaining, 30))
                    except asyncio.TimeoutError:
                        # Could be a long LLM call — keep waiting if under
                        # total timeout
                        continue

                    if msg.type == aiohttp.WSMsgType.TEXT:
                        data = json.loads(msg.data)
                        recv_record = {
                            "ts": time.time(),
                            "direction": "recv",
                            "turn": turn_label,
                            "type": data.get("type"),
                            "size": len(msg.data),
                            "data": data,
                        }
                        transcript.append(recv_record)
                        turn_record["ws_messages_received"].append({
                            "ts": time.time(),
                            "type": data.get("type"),
                            "size": len(msg.data),
                            "preview": json.dumps(data)[:500],
                        })

                        if data.get("type") == "response":
                            turn_record["response_text"] = data.get("content", "")
                            got_response = True

                        elif data.get("type") == "stream":
                            # Streaming response — accumulate
                            if turn_record["response_text"] is None:
                                turn_record["response_text"] = ""
                            turn_record["response_text"] += data.get("content", "")

                        elif data.get("type") == "stats":
                            turn_record["stats"] = data.get("data")
                            got_stats = True

                        elif data.get("type") == "error":
                            turn_record["error"] = data.get("content")
                            break

                        elif data.get("type") == "system_stats":
                            # This comes after stats — we're done
                            if got_stats:
                                break

                    elif msg.type in (aiohttp.WSMsgType.CLOSED,
                                      aiohttp.WSMsgType.ERROR):
                        turn_record["error"] = f"WebSocket closed/error: {msg.type}"
                        break

                turn_record["elapsed_s"] = round(time.time() - t0, 2)

    except Exception as e:
        turn_record["error"] = f"Connection error: {e}"
        turn_record["elapsed_s"] = round(time.time() - turn_record["sent_ts"], 2)

    return turn_record


# ---------------------------------------------------------------------------
# Journald capture
# ---------------------------------------------------------------------------

def start_journald_capture(output_path):
    """Start journalctl --user -u jarvis-web.service in background."""
    proc = subprocess.Popen(
        ["journalctl", "--user", "-u", "jarvis-web.service",
         "--since", "now", "-f", "--no-pager", "-o", "short-iso"],
        stdout=open(output_path, "w"),
        stderr=subprocess.STDOUT,
        preexec_fn=os.setsid,
    )
    return proc


def stop_journald_capture(proc):
    """Stop the journalctl process."""
    try:
        os.killpg(os.getpgid(proc.pid), signal.SIGTERM)
        proc.wait(timeout=5)
    except Exception:
        try:
            proc.kill()
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Sentinel management
# ---------------------------------------------------------------------------

SENTINEL_PATH = "/tmp/.jarvis_debug_active"


def activate_sentinel(debug_log_path):
    """Create sentinel file to activate structured debug logging."""
    with open(SENTINEL_PATH, "w") as f:
        f.write(debug_log_path)
    print(f"  Sentinel activated → {debug_log_path}")


def deactivate_sentinel():
    """Remove sentinel file to deactivate debug logging."""
    try:
        os.remove(SENTINEL_PATH)
        print("  Sentinel deactivated")
    except FileNotFoundError:
        pass


# ---------------------------------------------------------------------------
# Summary report
# ---------------------------------------------------------------------------

def generate_summary(turns, debug_log_path, journald_path, pptx_results,
                     share_dir):
    """Generate a human-readable consolidated report."""
    lines = []
    lines.append("=" * 80)
    lines.append("PRESENTATION ENGINE TEST REPORT")
    lines.append(f"Generated: {datetime.now().isoformat()}")
    lines.append("=" * 80)

    # --- Turn results ---
    lines.append("\n" + "─" * 80)
    lines.append("TURN RESULTS")
    lines.append("─" * 80)

    for turn in turns:
        lines.append(f"\n  [{turn['turn']}]")
        lines.append(f"  Sent: \"{turn['sent']}\"")
        lines.append(f"  Time: {turn['elapsed_s']}s")
        if turn["error"]:
            lines.append(f"  *** ERROR: {turn['error']}")
        resp = turn.get("response_text") or "(no response)"
        lines.append(f"  Response: \"{resp[:300]}\"")
        lines.append(f"  WS messages received: {len(turn['ws_messages_received'])}")
        for wm in turn["ws_messages_received"]:
            lines.append(f"    - {wm['type']} ({wm['size']} bytes)")
        if turn["stats"]:
            lines.append(f"  Stats: {json.dumps(turn['stats'], indent=4)}")

    # --- V36 append check ---
    lines.append("\n" + "─" * 80)
    lines.append("V36 APPEND FIX VALIDATION")
    lines.append("─" * 80)

    # Check if Turn 3 created a new file or modified the existing one
    if len(turns) >= 3:
        t1_resp = turns[0].get("response_text") or ""
        t3_resp = turns[2].get("response_text") or ""
        # Extract filenames from responses
        t1_file = None
        t3_file = None
        for m in re.finditer(r'(\S+\.pptx)', t1_resp):
            t1_file = m.group(1)
        for m in re.finditer(r'(\S+\.pptx)', t3_resp):
            t3_file = m.group(1)

        lines.append(f"  Turn 1 file: {t1_file}")
        lines.append(f"  Turn 3 file: {t3_file}")
        if t1_file and t3_file:
            if t1_file == t3_file:
                lines.append("  RESULT: PASS — Turn 3 modified the same file (append worked)")
            else:
                lines.append("  RESULT: FAIL — Turn 3 created a NEW file instead of appending")
                lines.append(f"         Expected: {t1_file}")
                lines.append(f"         Got:      {t3_file}")
        else:
            lines.append("  RESULT: INCONCLUSIVE — could not extract filenames from responses")

    # --- Pipeline cache check ---
    lines.append("\n" + "─" * 80)
    lines.append("PIPELINE CACHE VALIDATION")
    lines.append("─" * 80)

    if os.path.exists(debug_log_path):
        cache_events = []
        edit_events = []
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if evt.get("event") == "pipeline_cache_stored":
                    cache_events.append(evt)
                if evt.get("event") == "edit_presentation_cache_check":
                    edit_events.append(evt)
                if evt.get("event") == "edit_fallthrough_to_create":
                    edit_events.append(evt)

        lines.append(f"  Pipeline cache stored events: {len(cache_events)}")
        for ce in cache_events:
            lines.append(f"    - {ce.get('filename')}: {ce.get('slide_count')} slides, "
                         f"types={ce.get('slide_types')}")

        lines.append(f"  Edit cache check events: {len(edit_events)}")
        for ee in edit_events:
            lines.append(f"    - event={ee.get('event')}, has_cache={ee.get('has_cache')}, "
                         f"expired={ee.get('cache_expired')}, "
                         f"reason={ee.get('reason', 'N/A')}")

        if any(e.get("event") == "edit_fallthrough_to_create" for e in edit_events):
            lines.append("  *** PROBLEM: edit_presentation fell through to create_presentation")
            lines.append("      This means the pipeline cache was missing or expired.")
    else:
        lines.append("  *** Debug log not found — sentinel may not have activated")

    # --- Two-phase generation check ---
    lines.append("\n" + "─" * 80)
    lines.append("TWO-PHASE GENERATION VALIDATION")
    lines.append("─" * 80)

    if os.path.exists(debug_log_path):
        phase_events = []
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                etype = evt.get("event") or evt.get("type", "")
                if any(k in etype for k in ["two_phase", "phase1", "phase2"]):
                    phase_events.append(evt)

        lines.append(f"  Phase events captured: {len(phase_events)}")
        for pe in phase_events:
            etype = pe.get("event") or pe.get("type", "")
            lines.append(f"    - {etype}")
            if "layout" in pe:
                lines.append(f"      Layout: {pe['layout']}")
            if "type_names" in pe:
                lines.append(f"      Types: {pe['type_names']}")
            if "slide_count" in pe:
                lines.append(f"      Slides: {pe['slide_count']}")
            if "slide_types" in pe:
                lines.append(f"      Types: {pe['slide_types']}")
            if "raw_preview" in pe:
                preview = pe["raw_preview"][:300]
                lines.append(f"      LLM raw preview: {preview}")
            if "temperature" in pe:
                lines.append(f"      Temperature: {pe['temperature']}")

        if not phase_events:
            lines.append("  *** NO phase events found — two-phase generation may not have fired")
    else:
        lines.append("  *** Debug log not found")

    # --- Full debug log event summary ---
    lines.append("\n" + "─" * 80)
    lines.append("DEBUG LOG EVENT TIMELINE")
    lines.append("─" * 80)

    if os.path.exists(debug_log_path):
        event_counts = {}
        all_events = []
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                all_events.append(evt)
                key = evt.get("event") or evt.get("type", "unknown")
                event_counts[key] = event_counts.get(key, 0) + 1

        lines.append(f"  Total events: {len(all_events)}")
        lines.append(f"  Event type counts:")
        for k, v in sorted(event_counts.items()):
            lines.append(f"    {k}: {v}")

        lines.append(f"\n  Full timeline:")
        for evt in all_events:
            ts = evt.get("ts", 0)
            iso = datetime.fromtimestamp(ts).strftime("%H:%M:%S.%f")[:-3]
            etype = evt.get("event") or evt.get("type", "?")
            # Compact one-line summary per event
            detail_parts = []
            for dk in ["command", "skill_name", "intent", "layer",
                        "slide_count", "filename", "temperature",
                        "has_cache", "cache_expired", "reason",
                        "raw_len", "edit_request"]:
                if dk in evt:
                    val = evt[dk]
                    if isinstance(val, str) and len(val) > 80:
                        val = val[:80] + "..."
                    detail_parts.append(f"{dk}={val}")
            detail = ", ".join(detail_parts) if detail_parts else ""
            lines.append(f"    {iso} | {etype:40s} | {detail}")
    else:
        lines.append("  *** Debug log not found")

    # --- Routing check ---
    lines.append("\n" + "─" * 80)
    lines.append("ROUTING DECISIONS")
    lines.append("─" * 80)

    if os.path.exists(debug_log_path):
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if evt.get("type") == "route_decision":
                    cmd = evt.get("command", "")[:80]
                    mi = evt.get("match_info", {})
                    lines.append(f"  Command: \"{cmd}\"")
                    lines.append(f"    Handled: {evt.get('handled')}")
                    lines.append(f"    Skill: {mi.get('skill_name')}")
                    lines.append(f"    Handler: {mi.get('handler')}")
                    lines.append(f"    Confidence: {mi.get('confidence')}")
                    lines.append(f"    Routing time: {evt.get('routing_time_ms')}ms")
                    lines.append("")

    # --- PPTX inspection ---
    lines.append("\n" + "─" * 80)
    lines.append("PPTX FORMATTING INSPECTION")
    lines.append("─" * 80)

    for name, (result, text) in pptx_results.items():
        lines.append(f"\n  {name}:")
        if "error" in result:
            lines.append(f"    {result['error']}")
        else:
            lines.append(f"    Slides: {result['slide_count']}, "
                         f"Shapes: {result['total_shapes']}")
            lines.append(f"    Overflows: {result['overflows']}, "
                         f"Overlaps: {result['overlaps']}")
            lines.append(f"    Fonts: {result['fonts']}")
            lines.append(f"    Font sizes: {result['font_sizes']}")
            if result["issues"]:
                lines.append(f"    Issues ({len(result['issues'])}):")
                for iss in result["issues"]:
                    lines.append(f"      - {iss}")

    # --- Generated files check ---
    lines.append("\n" + "─" * 80)
    lines.append("GENERATED FILES IN SHARE/")
    lines.append("─" * 80)

    for f in sorted(Path(share_dir).glob("*.pptx")):
        mtime = datetime.fromtimestamp(f.stat().st_mtime).isoformat()
        lines.append(f"  {f.name:40s}  {f.stat().st_size:>10,} bytes  {mtime}")

    # --- Journald log size ---
    lines.append("\n" + "─" * 80)
    lines.append("JOURNALD CAPTURE")
    lines.append("─" * 80)
    if os.path.exists(journald_path):
        jsize = os.path.getsize(journald_path)
        with open(journald_path) as f:
            jlines = f.readlines()
        lines.append(f"  File: {journald_path}")
        lines.append(f"  Size: {jsize:,} bytes, {len(jlines)} lines")
        # Show key lines (errors, warnings, file_editor messages)
        key_lines = [l for l in jlines if any(
            k in l for k in ["ERROR", "WARNING", "file_editor",
                             "Pipeline cache", "Phase 1", "Phase 2",
                             "edit_presentation", "create_presentation",
                             "PPTX", "pptx"])]
        if key_lines:
            lines.append(f"  Key log lines ({len(key_lines)}):")
            for kl in key_lines:
                lines.append(f"    {kl.rstrip()}")
        else:
            lines.append("  No key log lines found (check full log)")
    else:
        lines.append(f"  *** Journald log not found at {journald_path}")

    lines.append("\n" + "=" * 80)
    lines.append("END OF REPORT")
    lines.append("=" * 80)

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

async def run_test(output_dir):
    """Run the full 3-turn test with all capture streams."""
    import aiohttp  # verify import early

    os.makedirs(output_dir, exist_ok=True)
    share_dir = os.path.expanduser("~/jarvis/share")

    # Determine run number
    existing = list(Path(output_dir).glob("pres_engine_*_summary.txt"))
    run_nums = []
    for f in existing:
        m = re.search(r'pres_engine_(\d+)_', f.name)
        if m:
            run_nums.append(int(m.group(1)))
    run_num = max(run_nums, default=0) + 1
    prefix = f"pres_engine_{run_num:03d}"

    ws_transcript_path = os.path.join(output_dir, f"{prefix}_ws_transcript.jsonl")
    debug_log_path = os.path.join(output_dir, f"{prefix}_debug.jsonl")
    journald_path = os.path.join(output_dir, f"{prefix}_journald.txt")
    pptx_inspect_path = os.path.join(output_dir, f"{prefix}_pptx_inspect.txt")
    summary_path = os.path.join(output_dir, f"{prefix}_summary.txt")

    print(f"\n{'='*60}")
    print(f"  Presentation Engine Test — Run {run_num:03d}")
    print(f"{'='*60}")
    print(f"  Output dir: {output_dir}")
    print(f"  Prefix: {prefix}")

    # --- Step 1: Clean old test presentations from share/ ---
    print(f"\n[Step 1] Cleaning old test presentations from share/...")
    for old in Path(share_dir).glob("*.pptx"):
        # Only remove files that look like test artifacts
        if old.name in ("ai_history_presentation.pptx",
                         "neural_networks_overview.pptx"):
            old.unlink()
            print(f"  Removed: {old.name}")

    # --- Step 2: Activate sentinel ---
    print(f"\n[Step 2] Activating debug sentinel...")
    activate_sentinel(debug_log_path)

    # --- Step 3: Start journald capture ---
    print(f"\n[Step 3] Starting journald capture...")
    journal_proc = start_journald_capture(journald_path)
    # Give journald a moment to start
    await asyncio.sleep(1)
    print(f"  Journald PID: {journal_proc.pid}")

    # --- Step 4: Run test turns ---
    ws_uri = "ws://localhost:8088/ws?token=QIj5fqyQ91HZ_SzAcOFJTaPB7Lu3Z570pxcrXIE944w"

    transcript = []  # shared across turns for ws_transcript output
    turns = []

    test_turns = [
        ("Turn 1: Create",
         "Create a 5 slide presentation about the history of artificial intelligence"),
        ("Turn 2: Edit",
         "Make slide 2 sound more technical"),
        ("Turn 3: Append",
         "Add a slide about neural networks"),
    ]

    for turn_label, message in test_turns:
        print(f"\n[Step 4] {turn_label}")
        print(f"  Sending: \"{message}\"")
        result = await ws_send_and_collect(
            ws_uri, message, turn_label, transcript, timeout=300)
        turns.append(result)
        print(f"  Elapsed: {result['elapsed_s']}s")
        if result["error"]:
            print(f"  *** ERROR: {result['error']}")
        else:
            resp_preview = (result.get("response_text") or "")[:120]
            print(f"  Response: \"{resp_preview}\"")
        print(f"  WS messages: {len(result['ws_messages_received'])}")

        # Brief pause between turns to let the server settle
        if turn_label != test_turns[-1][0]:
            await asyncio.sleep(2)

    # --- Step 5: Deactivate sentinel ---
    print(f"\n[Step 5] Deactivating sentinel...")
    deactivate_sentinel()

    # --- Step 6: Stop journald capture ---
    print(f"\n[Step 6] Stopping journald capture...")
    # Give a moment for final log lines to flush
    await asyncio.sleep(2)
    stop_journald_capture(journal_proc)
    if os.path.exists(journald_path):
        jsize = os.path.getsize(journald_path)
        print(f"  Captured: {jsize:,} bytes")

    # --- Step 7: Save WS transcript ---
    print(f"\n[Step 7] Saving WS transcript...")
    with open(ws_transcript_path, "w") as f:
        for entry in transcript:
            f.write(json.dumps(entry, default=str) + "\n")
    print(f"  Saved: {ws_transcript_path}")

    # --- Step 8: Inspect generated PPTX files ---
    print(f"\n[Step 8] Inspecting generated PPTX files...")
    pptx_results = {}
    inspect_text_parts = []
    for pptx_file in sorted(Path(share_dir).glob("*.pptx")):
        # Only inspect files modified after our test started
        result, text = inspect_pptx(str(pptx_file))
        pptx_results[pptx_file.name] = (result, text)
        inspect_text_parts.append(text)
        if isinstance(result, dict) and "error" not in result:
            print(f"  {pptx_file.name}: {result['slide_count']} slides, "
                  f"{result['overflows']} overflows, {result['overlaps']} overlaps")

    with open(pptx_inspect_path, "w") as f:
        f.write("\n\n".join(inspect_text_parts))
    print(f"  Saved: {pptx_inspect_path}")

    # --- Step 9: Generate summary report ---
    print(f"\n[Step 9] Generating summary report...")
    summary = generate_summary(turns, debug_log_path, journald_path,
                               pptx_results, share_dir)
    with open(summary_path, "w") as f:
        f.write(summary)
    print(f"  Saved: {summary_path}")

    # --- Done ---
    print(f"\n{'='*60}")
    print(f"  Test complete — Run {run_num:03d}")
    print(f"{'='*60}")
    print(f"\n  Output files:")
    print(f"    Summary:       {summary_path}")
    print(f"    WS transcript: {ws_transcript_path}")
    print(f"    Debug JSONL:   {debug_log_path}")
    print(f"    Journald:      {journald_path}")
    print(f"    PPTX inspect:  {pptx_inspect_path}")

    # Quick pass/fail summary
    print(f"\n  Quick results:")
    for turn in turns:
        status = "PASS" if not turn["error"] and turn["response_text"] else "FAIL"
        print(f"    {turn['turn']}: {status} ({turn['elapsed_s']}s)")

    # V36 check
    if len(turns) >= 3:
        t1_resp = turns[0].get("response_text") or ""
        t3_resp = turns[2].get("response_text") or ""
        t1_files = re.findall(r'(\S+\.pptx)', t1_resp)
        t3_files = re.findall(r'(\S+\.pptx)', t3_resp)
        if t1_files and t3_files:
            if t1_files[0] == t3_files[0]:
                print(f"    V36 Append: PASS (same file: {t1_files[0]})")
            else:
                print(f"    V36 Append: FAIL (Turn 1: {t1_files[0]}, Turn 3: {t3_files[0]})")

    return summary_path


def main():
    parser = argparse.ArgumentParser(
        description="Presentation Engine Test — Full Pipeline Diagnostics")
    parser.add_argument(
        "--output-dir",
        default=os.path.expanduser("~/jarvis/tests/iterative_results"),
        help="Directory for test output files")
    args = parser.parse_args()

    # Verify jarvis-web is running
    result = subprocess.run(
        ["systemctl", "--user", "is-active", "jarvis-web.service"],
        capture_output=True, text=True)
    if result.stdout.strip() != "active":
        print("ERROR: jarvis-web.service is not active.")
        print("Start it with: systemctl --user restart jarvis-web.service")
        sys.exit(1)

    # Verify aiohttp is available
    try:
        import aiohttp  # noqa: F401
    except ImportError:
        print("ERROR: aiohttp not installed. Run: pip install aiohttp")
        sys.exit(1)

    try:
        summary_path = asyncio.run(run_test(args.output_dir))
        print(f"\nDone. Read the full report:\n  cat {summary_path}")
    except KeyboardInterrupt:
        print("\nInterrupted — cleaning up sentinel...")
        deactivate_sentinel()
        sys.exit(130)
    finally:
        # Always clean up sentinel on exit
        deactivate_sentinel()


if __name__ == "__main__":
    main()
