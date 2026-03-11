#!/usr/bin/env python3
"""
Presentation Engine FLEX Test — Multi-Presentation Stress Test

Tests the full presentation pipeline across 3 wildly different topics,
with create/edit/remove/append operations and explicit theme requests.

Captures THREE data streams simultaneously:
  1. Client-side: every WebSocket message sent/received with timestamps
  2. Server-side: journald logs from jarvis-web.service
  3. Debug JSONL: structured pipeline events via sentinel file

Sequences:
  A. Ancient Rome (12 slides, elegant theme) — create, edit, remove, add
  B. Cybersecurity Threat Hunting (10 slides, jarvis theme) — create, edit, remove, add
  C. Science of Cooking (12 slides, earth theme) — create, edit, remove, add

Usage:
    python3 scripts/test_pres_flex.py [--output-dir DIR]

Outputs (all in output dir):
    pres_flex_NNN_ws_transcript.jsonl   — raw WebSocket messages
    pres_flex_NNN_debug.jsonl           — structured debug pipeline log
    pres_flex_NNN_journald.txt          — journald log capture
    pres_flex_NNN_pptx_inspect.txt      — PPTX formatting inspection
    pres_flex_NNN_summary.txt           — human-readable consolidated report
"""

import argparse
import asyncio
import json
import os
import re
import signal
import subprocess
import sys
import time
from datetime import datetime, timezone
from pathlib import Path

# ---------------------------------------------------------------------------
# PPTX inspection (same as test_pres_engine.py)
# ---------------------------------------------------------------------------

EMU_PER_INCH = 914400


def _emu_to_in(emu):
    if emu is None:
        return 0.0
    return round(emu / EMU_PER_INCH, 3)


def _estimate_text_overflow(text_frame, shape):
    """Return (needed_inches, shape_inches, overflows_bool)."""
    total_lines = 0
    max_fs = 12
    for para in text_frame.paragraphs:
        text = para.text
        if not text.strip():
            total_lines += 0.5
            continue
        fs = None
        for run in para.runs:
            if run.font.size:
                fs = run.font.size / 12700
                break
        if fs is None:
            fs = max_fs
        else:
            max_fs = max(max_fs, fs)
        shape_w_pt = (shape.width / EMU_PER_INCH) * 72
        chars_per_line = max(1, shape_w_pt / (0.5 * fs))
        total_lines += max(1, len(text) / chars_per_line)
    needed = (total_lines * max_fs * 1.4) / 72
    actual = shape.height / EMU_PER_INCH
    return round(needed, 2), round(actual, 2), needed > actual * 1.1


def _check_overlaps(shapes_info):
    """Return list of (shape_a, shape_b, overlap_sq_in).

    Skips containment pairs (one shape fully inside another) since those
    are intentional (e.g. TextBox inside a RoundedRectangle card).
    """
    overlaps = []
    for i, s1 in enumerate(shapes_info):
        for j, s2 in enumerate(shapes_info):
            if j <= i:
                continue
            l1, t1, r1, b1 = s1
            l2, t2, r2, b2 = s2
            if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                # Skip containment (one shape fully inside the other)
                if (l1 <= l2 and t1 <= t2 and r1 >= r2 and b1 >= b2) or \
                   (l2 <= l1 and t2 <= t1 and r2 >= r1 and b2 >= b1):
                    continue
                area = ((min(r1, r2) - max(l1, l2)) *
                        (min(b1, b2) - max(t1, t2))) / (EMU_PER_INCH ** 2)
                overlaps.append((i, j, round(area, 3)))
    return overlaps


def inspect_pptx(filepath):
    """Inspect a PPTX file and return structured results + human text."""
    from pptx import Presentation as PptxPres

    lines = []
    issues = []

    if not os.path.exists(filepath):
        return {"error": f"File not found: {filepath}"}, f"FILE NOT FOUND: {filepath}\n"

    prs = PptxPres(filepath)
    sw, sh = prs.slide_width, prs.slide_height
    lines.append(f"FILE: {os.path.basename(filepath)}")
    lines.append(f"  Size on disk: {os.path.getsize(filepath):,} bytes")
    lines.append(f"  Slide dimensions: {_emu_to_in(sw)}\" x {_emu_to_in(sh)}\"")
    lines.append(f"  Slide count: {len(prs.slides)}")

    all_fonts = set()
    all_sizes = set()
    total_overflows = 0
    total_overlaps = 0
    slide_details = []

    for si, slide in enumerate(prs.slides):
        snum = si + 1
        lines.append(f"\n  --- Slide {snum} ---")
        try:
            lines.append(f"  Layout: {slide.slide_layout.name}")
        except Exception:
            lines.append("  Layout: unknown")

        shapes_bounds = []
        shape_names = []
        slide_overflows = 0
        slide_overlap_count = 0

        for shape in slide.shapes:
            left = shape.left or 0
            top = shape.top or 0
            width = shape.width or 0
            height = shape.height or 0
            shapes_bounds.append((left, top, left + width, top + height))
            shape_names.append(shape.name or "unnamed")

            oob = []
            if left < 0:
                oob.append("left<0")
            if top < 0:
                oob.append("top<0")
            if left + width > sw:
                oob.append(f"right={_emu_to_in(left+width)}\">{_emu_to_in(sw)}\"")
            if top + height > sh:
                oob.append(f"bottom={_emu_to_in(top+height)}\">{_emu_to_in(sh)}\"")

            shape_desc = (f"    [{shape.name}] type={shape.shape_type} "
                          f"pos=({_emu_to_in(left)}\",{_emu_to_in(top)}\") "
                          f"size=({_emu_to_in(width)}\"x{_emu_to_in(height)}\")")
            if oob:
                shape_desc += f" *** OOB: {', '.join(oob)}"
                issues.append(f"Slide {snum}: {shape.name} out of bounds — {', '.join(oob)}")

            if shape.has_text_frame:
                tf = shape.text_frame
                needed, actual, overflows = _estimate_text_overflow(tf, shape)
                text_preview = tf.text[:80].replace('\n', ' ')
                if overflows and actual > 0.1:
                    slide_overflows += 1
                    total_overflows += 1
                    shape_desc += f" *** OVERFLOW: needs ~{needed}\" has {actual}\""
                    issues.append(
                        f"Slide {snum}: {shape.name} text overflow — "
                        f"needs ~{needed}\" in {actual}\" container "
                        f"(\"{text_preview[:60]}...\")")

                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            all_fonts.add(run.font.name)
                        if run.font.size:
                            all_sizes.add(round(run.font.size / 12700, 1))

            lines.append(shape_desc)

        overlaps = _check_overlaps(shapes_bounds)
        if overlaps:
            for a, b, area in overlaps:
                if area > 0.5:
                    total_overlaps += 1
                    slide_overlap_count += 1
                    issues.append(
                        f"Slide {snum}: {shape_names[a]} <-> {shape_names[b]} "
                        f"overlap {area} sq in")
            if slide_overlap_count:
                lines.append(f"    *** {slide_overlap_count} significant overlap(s) on this slide")

        slide_details.append({
            "slide_num": snum,
            "shape_count": len(shapes_bounds),
            "overflows": slide_overflows,
            "overlaps": slide_overlap_count,
        })

    lines.append(f"\n  SUMMARY: {len(prs.slides)} slides, "
                 f"{sum(s['shape_count'] for s in slide_details)} shapes")
    lines.append(f"  Fonts: {sorted(all_fonts) if all_fonts else 'inherited'}")
    lines.append(f"  Font sizes: {sorted(all_sizes) if all_sizes else 'inherited'}")
    lines.append(f"  Text overflows: {total_overflows}")
    lines.append(f"  Significant overlaps: {total_overlaps}")
    if issues:
        lines.append(f"\n  ALL ISSUES ({len(issues)}):")
        for iss in issues:
            lines.append(f"    - {iss}")

    result = {
        "file": os.path.basename(filepath),
        "slide_count": len(prs.slides),
        "total_shapes": sum(s['shape_count'] for s in slide_details),
        "fonts": sorted(all_fonts),
        "font_sizes": sorted(all_sizes),
        "overflows": total_overflows,
        "overlaps": total_overlaps,
        "issues": issues,
        "slides": slide_details,
    }
    return result, "\n".join(lines)


# ---------------------------------------------------------------------------
# WebSocket client
# ---------------------------------------------------------------------------

async def ws_drain_initial(ws):
    """Drain initial burst messages after connecting (until system_stats)."""
    import aiohttp

    initial_messages = []
    try:
        while True:
            msg = await asyncio.wait_for(ws.receive(), timeout=5)
            if msg.type == aiohttp.WSMsgType.TEXT:
                data = json.loads(msg.data)
                initial_messages.append({
                    "ts": time.time(),
                    "direction": "recv_initial",
                    "type": data.get("type"),
                    "size": len(msg.data),
                })
                if data.get("type") == "system_stats":
                    break
            elif msg.type in (aiohttp.WSMsgType.CLOSED,
                              aiohttp.WSMsgType.ERROR):
                break
    except asyncio.TimeoutError:
        pass
    return initial_messages


async def ws_send_and_collect(ws, message_text, turn_label, transcript,
                              timeout=600):
    """Send a chat message on an existing WS and collect all response messages."""
    import aiohttp

    turn_record = {
        "turn": turn_label,
        "sent": message_text,
        "sent_ts": time.time(),
        "sent_iso": datetime.now(timezone.utc).isoformat(),
        "ws_messages_received": [],
        "response_text": None,
        "stats": None,
        "error": None,
        "elapsed_s": None,
    }

    try:
        # Send message
        payload = json.dumps({"type": "message", "content": message_text})
        await ws.send_str(payload)
        transcript.append({
            "ts": time.time(),
            "direction": "sent",
            "turn": turn_label,
            "payload": {"type": "message", "content": message_text},
        })

        t0 = time.time()

        # Collect responses until stats message (end-of-response)
        got_stats = False
        while True:
            remaining = timeout - (time.time() - t0)
            if remaining <= 0:
                turn_record["error"] = f"Timeout after {timeout}s"
                break
            try:
                msg = await asyncio.wait_for(
                    ws.receive(),
                    timeout=min(remaining, 60))
            except asyncio.TimeoutError:
                continue

            if msg.type == aiohttp.WSMsgType.TEXT:
                data = json.loads(msg.data)
                recv_record = {
                    "ts": time.time(),
                    "direction": "recv",
                    "turn": turn_label,
                    "type": data.get("type"),
                    "size": len(msg.data),
                    "data": data,
                }
                transcript.append(recv_record)
                turn_record["ws_messages_received"].append({
                    "ts": time.time(),
                    "type": data.get("type"),
                    "size": len(msg.data),
                    "preview": json.dumps(data)[:500],
                })

                if data.get("type") == "response":
                    turn_record["response_text"] = data.get("content", "")

                elif data.get("type") == "stream":
                    if turn_record["response_text"] is None:
                        turn_record["response_text"] = ""
                    turn_record["response_text"] += data.get("content", "")

                elif data.get("type") == "stream_token":
                    if turn_record["response_text"] is None:
                        turn_record["response_text"] = ""
                    turn_record["response_text"] += data.get("token", "")

                elif data.get("type") == "stream_end":
                    pass  # response text already accumulated from stream_token

                elif data.get("type") == "stats":
                    turn_record["stats"] = data.get("data")
                    got_stats = True

                elif data.get("type") == "error":
                    turn_record["error"] = data.get("content")
                    break

                elif data.get("type") == "system_stats":
                    if got_stats:
                        break

            elif msg.type in (aiohttp.WSMsgType.CLOSED,
                              aiohttp.WSMsgType.ERROR):
                turn_record["error"] = f"WebSocket closed/error: {msg.type}"
                break

        turn_record["elapsed_s"] = round(time.time() - t0, 2)

    except Exception as e:
        turn_record["error"] = f"Connection error: {e}"
        turn_record["elapsed_s"] = round(time.time() - turn_record["sent_ts"], 2)

    return turn_record


# ---------------------------------------------------------------------------
# Journald capture
# ---------------------------------------------------------------------------

def start_journald_capture(output_path):
    proc = subprocess.Popen(
        ["journalctl", "--user", "-u", "jarvis-web.service",
         "--since", "now", "-f", "--no-pager", "-o", "short-iso"],
        stdout=open(output_path, "w"),
        stderr=subprocess.STDOUT,
        preexec_fn=os.setsid,
    )
    return proc


def stop_journald_capture(proc):
    try:
        os.killpg(os.getpgid(proc.pid), signal.SIGTERM)
        proc.wait(timeout=5)
    except Exception:
        try:
            proc.kill()
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Sentinel management
# ---------------------------------------------------------------------------

SENTINEL_PATH = "/tmp/.jarvis_debug_active"


def activate_sentinel(debug_log_path):
    with open(SENTINEL_PATH, "w") as f:
        f.write(debug_log_path)
    print(f"  Sentinel activated -> {debug_log_path}")


def deactivate_sentinel():
    try:
        os.remove(SENTINEL_PATH)
        print("  Sentinel deactivated")
    except FileNotFoundError:
        pass


# ---------------------------------------------------------------------------
# Test sequences
# ---------------------------------------------------------------------------

SEQUENCES = [
    {
        "name": "A: Ancient Rome",
        "theme": "elegant",
        "turns": [
            ("A1-Create",
             "Create a 12 slide presentation about the rise and fall of "
             "Ancient Rome. Use the elegant theme."),
            ("A2-Edit",
             "Make slide 3 more dramatic and vivid — add imagery of Roman "
             "legions and conquest."),
            ("A3-Remove",
             "Remove slide 5 from the presentation."),
            ("A4-Append",
             "Add a slide about Roman engineering marvels — aqueducts, roads, "
             "and the Colosseum."),
        ],
    },
    {
        "name": "B: Cybersecurity",
        "theme": "jarvis",
        "turns": [
            ("B1-Create",
             "Create a 10 slide presentation about cybersecurity threat "
             "hunting techniques. Use the jarvis theme."),
            ("B2-Edit",
             "Make the introduction slide more engaging and less dry — add a "
             "compelling hook about recent breaches."),
            ("B3-Remove",
             "Remove slide 7 from the presentation."),
            ("B4-Append",
             "Add two slides about incident response best practices and "
             "post-incident review."),
        ],
    },
    {
        "name": "C: Cooking Science",
        "theme": "earth",
        "turns": [
            ("C1-Create",
             "Create a 12 slide presentation about the science of cooking "
             "and kitchen chemistry. Use the earth theme."),
            ("C2-Edit",
             "Make slide 6 more technical with specific chemical reactions "
             "like the Maillard reaction and caramelization."),
            ("C3-Remove",
             "Remove slide 2 from the presentation."),
            ("C4-Append",
             "Add a slide about molecular gastronomy techniques like "
             "spherification and sous vide."),
        ],
    },
]


# ---------------------------------------------------------------------------
# Summary report
# ---------------------------------------------------------------------------

def generate_summary(sequences_results, debug_log_path, journald_path,
                     pptx_results, share_dir):
    """Generate human-readable consolidated report."""
    lines = []
    lines.append("=" * 80)
    lines.append("PRESENTATION ENGINE FLEX TEST REPORT")
    lines.append(f"Generated: {datetime.now().isoformat()}")
    lines.append(f"Sequences: {len(SEQUENCES)}, Total turns: "
                 f"{sum(len(s['results']) for s in sequences_results)}")
    lines.append("=" * 80)

    # --- Sequence results ---
    for seq in sequences_results:
        lines.append(f"\n{'─'*80}")
        lines.append(f"SEQUENCE: {seq['name']} (theme: {seq['theme']})")
        lines.append(f"{'─'*80}")

        for turn in seq["results"]:
            status = "PASS" if not turn["error"] and turn["response_text"] else "FAIL"
            lines.append(f"\n  [{turn['turn']}] — {status}")
            lines.append(f"  Sent: \"{turn['sent']}\"")
            lines.append(f"  Time: {turn['elapsed_s']}s")
            if turn["error"]:
                lines.append(f"  *** ERROR: {turn['error']}")
            resp = turn.get("response_text") or "(no response)"
            # Show more of the response for flex test
            lines.append(f"  Response: \"{resp[:500]}\"")
            lines.append(f"  WS messages received: {len(turn['ws_messages_received'])}")
            for wm in turn["ws_messages_received"]:
                lines.append(f"    - {wm['type']} ({wm['size']} bytes)")
            if turn["stats"]:
                lines.append(f"  Stats: {json.dumps(turn['stats'], indent=4)}")

    # --- Per-sequence file tracking ---
    lines.append(f"\n{'─'*80}")
    lines.append("FILE TRACKING PER SEQUENCE")
    lines.append(f"{'─'*80}")

    for seq in sequences_results:
        lines.append(f"\n  {seq['name']}:")
        files_mentioned = set()
        for turn in seq["results"]:
            resp = turn.get("response_text") or ""
            for m in re.finditer(r'(\S+\.pptx)', resp):
                files_mentioned.add(m.group(1))
        if files_mentioned:
            for fn in sorted(files_mentioned):
                lines.append(f"    - {fn}")
        else:
            lines.append("    (no .pptx filenames found in responses)")

        # Check create vs append consistency
        create_turn = seq["results"][0]
        append_turn = seq["results"][-1]  # last turn is always an append
        create_resp = create_turn.get("response_text") or ""
        append_resp = append_turn.get("response_text") or ""
        create_files = re.findall(r'(\S+\.pptx)', create_resp)
        append_files = re.findall(r'(\S+\.pptx)', append_resp)
        if create_files and append_files:
            if create_files[0] == append_files[0]:
                lines.append(f"    Append check: PASS (same file: {create_files[0]})")
            else:
                lines.append(f"    Append check: FAIL — create={create_files[0]}, "
                             f"append={append_files[0]}")
        else:
            lines.append("    Append check: INCONCLUSIVE")

    # --- Pipeline cache analysis ---
    lines.append(f"\n{'─'*80}")
    lines.append("PIPELINE CACHE ANALYSIS")
    lines.append(f"{'─'*80}")

    if os.path.exists(debug_log_path):
        cache_events = []
        edit_events = []
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if evt.get("event") == "pipeline_cache_stored":
                    cache_events.append(evt)
                if "edit" in (evt.get("event") or "").lower() or \
                   "cache_check" in (evt.get("event") or "").lower():
                    edit_events.append(evt)
                if evt.get("event") == "edit_fallthrough_to_create":
                    edit_events.append(evt)

        lines.append(f"  Cache stored events: {len(cache_events)}")
        for ce in cache_events:
            lines.append(f"    - {ce.get('filename')}: {ce.get('slide_count')} slides, "
                         f"types={ce.get('slide_types')}")
        lines.append(f"  Edit/cache events: {len(edit_events)}")
        for ee in edit_events:
            lines.append(f"    - {ee.get('event')}: "
                         f"has_cache={ee.get('has_cache')}, "
                         f"expired={ee.get('cache_expired')}, "
                         f"reason={ee.get('reason', 'N/A')}")
    else:
        lines.append("  *** Debug log not found")

    # --- Two-phase generation ---
    lines.append(f"\n{'─'*80}")
    lines.append("TWO-PHASE GENERATION")
    lines.append(f"{'─'*80}")

    if os.path.exists(debug_log_path):
        phase_events = []
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                etype = evt.get("event") or evt.get("type", "")
                if any(k in etype for k in ["two_phase", "phase1", "phase2"]):
                    phase_events.append(evt)

        lines.append(f"  Phase events: {len(phase_events)}")
        for pe in phase_events:
            etype = pe.get("event") or pe.get("type", "")
            detail_parts = []
            for k in ["layout", "type_names", "slide_count", "slide_types",
                       "temperature", "raw_preview"]:
                if k in pe:
                    val = pe[k]
                    if isinstance(val, str) and len(val) > 200:
                        val = val[:200] + "..."
                    detail_parts.append(f"{k}={val}")
            lines.append(f"    - {etype}: {', '.join(detail_parts)}")
        if not phase_events:
            lines.append("  *** No phase events found")
    else:
        lines.append("  *** Debug log not found")

    # --- Full debug timeline ---
    lines.append(f"\n{'─'*80}")
    lines.append("DEBUG LOG EVENT TIMELINE")
    lines.append(f"{'─'*80}")

    if os.path.exists(debug_log_path):
        event_counts = {}
        all_events = []
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                all_events.append(evt)
                key = evt.get("event") or evt.get("type", "unknown")
                event_counts[key] = event_counts.get(key, 0) + 1

        lines.append(f"  Total events: {len(all_events)}")
        lines.append(f"  Event type counts:")
        for k, v in sorted(event_counts.items()):
            lines.append(f"    {k}: {v}")

        lines.append(f"\n  Full timeline:")
        for evt in all_events:
            ts = evt.get("ts", 0)
            try:
                iso = datetime.fromtimestamp(ts).strftime("%H:%M:%S.%f")[:-3]
            except (OSError, ValueError):
                iso = "??:??:??.???"
            etype = evt.get("event") or evt.get("type", "?")
            detail_parts = []
            for dk in ["command", "skill_name", "intent", "layer",
                        "slide_count", "filename", "temperature",
                        "has_cache", "cache_expired", "reason",
                        "raw_len", "edit_request"]:
                if dk in evt:
                    val = evt[dk]
                    if isinstance(val, str) and len(val) > 80:
                        val = val[:80] + "..."
                    detail_parts.append(f"{dk}={val}")
            detail = ", ".join(detail_parts) if detail_parts else ""
            lines.append(f"    {iso} | {etype:40s} | {detail}")
    else:
        lines.append("  *** Debug log not found")

    # --- Routing decisions ---
    lines.append(f"\n{'─'*80}")
    lines.append("ROUTING DECISIONS")
    lines.append(f"{'─'*80}")

    if os.path.exists(debug_log_path):
        with open(debug_log_path) as f:
            for line in f:
                try:
                    evt = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if evt.get("type") == "route_decision":
                    cmd = evt.get("command", "")[:100]
                    mi = evt.get("match_info", {})
                    lines.append(f"  Command: \"{cmd}\"")
                    lines.append(f"    Handled: {evt.get('handled')}")
                    lines.append(f"    Skill: {mi.get('skill_name')}")
                    lines.append(f"    Handler: {mi.get('handler')}")
                    lines.append(f"    Confidence: {mi.get('confidence')}")
                    lines.append(f"    Routing time: {evt.get('routing_time_ms')}ms")
                    lines.append("")

    # --- PPTX inspection ---
    lines.append(f"\n{'─'*80}")
    lines.append("PPTX FORMATTING INSPECTION")
    lines.append(f"{'─'*80}")

    for name, (result, text) in pptx_results.items():
        lines.append(f"\n  {name}:")
        if "error" in result:
            lines.append(f"    {result['error']}")
        else:
            lines.append(f"    Slides: {result['slide_count']}, "
                         f"Shapes: {result['total_shapes']}")
            lines.append(f"    Overflows: {result['overflows']}, "
                         f"Overlaps: {result['overlaps']}")
            lines.append(f"    Fonts: {result['fonts']}")
            lines.append(f"    Font sizes: {result['font_sizes']}")
            if result["issues"]:
                lines.append(f"    Issues ({len(result['issues'])}):")
                for iss in result["issues"]:
                    lines.append(f"      - {iss}")
            else:
                lines.append(f"    Issues: NONE")

    # --- Theme validation ---
    lines.append(f"\n{'─'*80}")
    lines.append("THEME VALIDATION")
    lines.append(f"{'─'*80}")

    # Expected fonts per theme
    expected_fonts = {
        "elegant": {"heading": "Book Antiqua", "body": "Verdana"},
        "jarvis": {"heading": "Century Gothic", "body": "Calibri"},
        "earth": {"heading": "Georgia", "body": "Verdana"},
        "professional": {"heading": "Cambria", "body": "Calibri"},
    }

    for name, (result, _) in pptx_results.items():
        if "error" in result:
            continue
        fonts_found = result.get("fonts", [])
        lines.append(f"\n  {name}:")
        lines.append(f"    Fonts found: {fonts_found}")

        # Try to match against expected themes
        matched_theme = None
        for theme, expected in expected_fonts.items():
            if expected["heading"] in fonts_found or expected["body"] in fonts_found:
                matched_theme = theme
                break
        if matched_theme:
            lines.append(f"    Likely theme: {matched_theme}")
            exp = expected_fonts[matched_theme]
            if exp["heading"] in fonts_found:
                lines.append(f"    Heading font ({exp['heading']}): PRESENT")
            else:
                lines.append(f"    Heading font ({exp['heading']}): MISSING")
            if exp["body"] in fonts_found:
                lines.append(f"    Body font ({exp['body']}): PRESENT")
            else:
                lines.append(f"    Body font ({exp['body']}): MISSING")
        else:
            lines.append(f"    Theme: COULD NOT DETERMINE (all fonts inherited?)")

    # --- Generated files ---
    lines.append(f"\n{'─'*80}")
    lines.append("GENERATED FILES IN SHARE/")
    lines.append(f"{'─'*80}")

    for f in sorted(Path(share_dir).glob("*.pptx")):
        mtime = datetime.fromtimestamp(f.stat().st_mtime).isoformat()
        lines.append(f"  {f.name:50s}  {f.stat().st_size:>10,} bytes  {mtime}")

    # --- Journald capture ---
    lines.append(f"\n{'─'*80}")
    lines.append("JOURNALD CAPTURE")
    lines.append(f"{'─'*80}")

    if os.path.exists(journald_path):
        jsize = os.path.getsize(journald_path)
        with open(journald_path) as f:
            jlines = f.readlines()
        lines.append(f"  File: {journald_path}")
        lines.append(f"  Size: {jsize:,} bytes, {len(jlines)} lines")
        key_lines = [l for l in jlines if any(
            k in l for k in ["ERROR", "WARNING", "file_editor",
                             "Pipeline cache", "Phase 1", "Phase 2",
                             "edit_presentation", "create_presentation",
                             "PPTX", "pptx", "theme", "accent"])]
        if key_lines:
            lines.append(f"  Key log lines ({len(key_lines)}):")
            for kl in key_lines[:100]:  # cap at 100 to avoid massive reports
                lines.append(f"    {kl.rstrip()}")
            if len(key_lines) > 100:
                lines.append(f"    ... ({len(key_lines) - 100} more)")
        else:
            lines.append("  No key log lines found (check full log)")
    else:
        lines.append(f"  *** Journald log not found at {journald_path}")

    # --- Overall scorecard ---
    lines.append(f"\n{'='*80}")
    lines.append("OVERALL SCORECARD")
    lines.append(f"{'='*80}")

    total_turns = sum(len(s["results"]) for s in sequences_results)
    pass_turns = sum(
        1 for s in sequences_results for t in s["results"]
        if not t["error"] and t["response_text"])
    fail_turns = total_turns - pass_turns

    total_files = len(pptx_results)
    total_slides = sum(
        r["slide_count"] for r, _ in pptx_results.values()
        if isinstance(r, dict) and "slide_count" in r)
    total_overflows = sum(
        r["overflows"] for r, _ in pptx_results.values()
        if isinstance(r, dict) and "overflows" in r)
    total_overlaps = sum(
        r["overlaps"] for r, _ in pptx_results.values()
        if isinstance(r, dict) and "overlaps" in r)

    lines.append(f"  Turns: {pass_turns}/{total_turns} PASS, {fail_turns} FAIL")
    lines.append(f"  Files generated: {total_files}")
    lines.append(f"  Total slides: {total_slides}")
    lines.append(f"  Total text overflows: {total_overflows}")
    lines.append(f"  Total shape overlaps: {total_overlaps}")

    total_time = sum(
        t["elapsed_s"] or 0 for s in sequences_results for t in s["results"])
    lines.append(f"  Total generation time: {total_time:.1f}s")

    lines.append(f"\n{'='*80}")
    lines.append("END OF REPORT")
    lines.append(f"{'='*80}")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

async def run_test(output_dir):
    """Run the full multi-sequence flex test."""
    import aiohttp  # noqa: F401 — verify import early

    os.makedirs(output_dir, exist_ok=True)
    share_dir = os.path.expanduser("~/jarvis/share")

    # Determine run number
    existing = list(Path(output_dir).glob("pres_flex_*_summary.txt"))
    run_nums = []
    for f in existing:
        m = re.search(r'pres_flex_(\d+)_', f.name)
        if m:
            run_nums.append(int(m.group(1)))
    run_num = max(run_nums, default=0) + 1
    prefix = f"pres_flex_{run_num:03d}"

    ws_transcript_path = os.path.join(output_dir, f"{prefix}_ws_transcript.jsonl")
    debug_log_path = os.path.join(output_dir, f"{prefix}_debug.jsonl")
    journald_path = os.path.join(output_dir, f"{prefix}_journald.txt")
    pptx_inspect_path = os.path.join(output_dir, f"{prefix}_pptx_inspect.txt")
    summary_path = os.path.join(output_dir, f"{prefix}_summary.txt")

    print(f"\n{'='*60}")
    print(f"  Presentation Engine FLEX Test — Run {run_num:03d}")
    print(f"  {len(SEQUENCES)} sequences, "
          f"{sum(len(s['turns']) for s in SEQUENCES)} turns")
    print(f"{'='*60}")
    print(f"  Output dir: {output_dir}")
    print(f"  Prefix: {prefix}")

    # --- Step 1: Activate sentinel ---
    print(f"\n[Step 1] Activating debug sentinel...")
    activate_sentinel(debug_log_path)

    # --- Step 2: Start journald capture ---
    print(f"\n[Step 2] Starting journald capture...")
    journal_proc = start_journald_capture(journald_path)
    await asyncio.sleep(1)
    print(f"  Journald PID: {journal_proc.pid}")

    # --- Step 3: Run sequences ---
    ws_uri = "ws://localhost:8088/ws?token=QIj5fqyQ91HZ_SzAcOFJTaPB7Lu3Z570pxcrXIE944w"
    transcript = []
    sequences_results = []

    import aiohttp

    for seq_idx, seq in enumerate(SEQUENCES):
        seq_result = {
            "name": seq["name"],
            "theme": seq["theme"],
            "results": [],
        }

        print(f"\n{'─'*60}")
        print(f"  SEQUENCE {seq['name']} (theme: {seq['theme']})")
        print(f"{'─'*60}")

        # Keep a SINGLE WebSocket connection open for all turns in this sequence
        # so the server preserves conversation context between turns.
        async with aiohttp.ClientSession() as session:
            async with session.ws_connect(ws_uri, timeout=30) as ws:
                initial = await ws_drain_initial(ws)
                print(f"  WS connected, drained {len(initial)} initial messages")

                for turn_idx, (turn_label, message) in enumerate(seq["turns"]):
                    print(f"\n  [{turn_label}]")
                    print(f"  Sending: \"{message[:100]}{'...' if len(message) > 100 else ''}\"")

                    result = await ws_send_and_collect(
                        ws, message, turn_label, transcript, timeout=600)
                    result["initial_messages"] = initial if turn_idx == 0 else []
                    seq_result["results"].append(result)

                    print(f"  Elapsed: {result['elapsed_s']}s")
                    if result["error"]:
                        print(f"  *** ERROR: {result['error']}")
                    else:
                        resp_preview = (result.get("response_text") or "")[:150]
                        print(f"  Response: \"{resp_preview}\"")
                    print(f"  WS messages: {len(result['ws_messages_received'])}")

                    # Pause between turns within a sequence
                    if turn_idx < len(seq["turns"]) - 1:
                        await asyncio.sleep(3)

        sequences_results.append(seq_result)

        # Longer pause between sequences to let server settle
        if seq_idx < len(SEQUENCES) - 1:
            print(f"\n  --- Pausing 5s before next sequence ---")
            await asyncio.sleep(5)

    # --- Step 4: Deactivate sentinel ---
    print(f"\n[Step 4] Deactivating sentinel...")
    deactivate_sentinel()

    # --- Step 5: Stop journald capture ---
    print(f"\n[Step 5] Stopping journald capture...")
    await asyncio.sleep(2)
    stop_journald_capture(journal_proc)
    if os.path.exists(journald_path):
        jsize = os.path.getsize(journald_path)
        print(f"  Captured: {jsize:,} bytes")

    # --- Step 6: Save WS transcript ---
    print(f"\n[Step 6] Saving WS transcript...")
    with open(ws_transcript_path, "w") as f:
        for entry in transcript:
            f.write(json.dumps(entry, default=str) + "\n")
    print(f"  Saved: {ws_transcript_path}")

    # --- Step 7: Inspect generated PPTX files ---
    print(f"\n[Step 7] Inspecting generated PPTX files...")
    pptx_results = {}
    inspect_text_parts = []
    for pptx_file in sorted(Path(share_dir).glob("*.pptx")):
        result, text = inspect_pptx(str(pptx_file))
        pptx_results[pptx_file.name] = (result, text)
        inspect_text_parts.append(text)
        if isinstance(result, dict) and "error" not in result:
            print(f"  {pptx_file.name}: {result['slide_count']} slides, "
                  f"{result['overflows']} overflows, {result['overlaps']} overlaps, "
                  f"fonts={result['fonts'][:3]}{'...' if len(result['fonts']) > 3 else ''}")

    with open(pptx_inspect_path, "w") as f:
        f.write("\n\n".join(inspect_text_parts))
    print(f"  Saved: {pptx_inspect_path}")

    # --- Step 8: Generate summary report ---
    print(f"\n[Step 8] Generating summary report...")
    summary = generate_summary(sequences_results, debug_log_path, journald_path,
                               pptx_results, share_dir)
    with open(summary_path, "w") as f:
        f.write(summary)
    print(f"  Saved: {summary_path}")

    # --- Done ---
    print(f"\n{'='*60}")
    print(f"  FLEX Test complete — Run {run_num:03d}")
    print(f"{'='*60}")
    print(f"\n  Output files:")
    print(f"    Summary:       {summary_path}")
    print(f"    WS transcript: {ws_transcript_path}")
    print(f"    Debug JSONL:   {debug_log_path}")
    print(f"    Journald:      {journald_path}")
    print(f"    PPTX inspect:  {pptx_inspect_path}")

    # Quick scorecard
    print(f"\n  SCORECARD:")
    for seq in sequences_results:
        print(f"\n  {seq['name']}:")
        for turn in seq["results"]:
            status = "PASS" if not turn["error"] and turn["response_text"] else "FAIL"
            elapsed = turn["elapsed_s"] or 0
            print(f"    {turn['turn']:15s}: {status} ({elapsed:.1f}s)")

    total_turns = sum(len(s["results"]) for s in sequences_results)
    pass_count = sum(
        1 for s in sequences_results for t in s["results"]
        if not t["error"] and t["response_text"])
    print(f"\n  Overall: {pass_count}/{total_turns} turns passed")

    return summary_path


def main():
    parser = argparse.ArgumentParser(
        description="Presentation Engine FLEX Test — Multi-Presentation Stress Test")
    parser.add_argument(
        "--output-dir",
        default=os.path.expanduser("~/jarvis/tests/iterative_results"),
        help="Directory for test output files")
    args = parser.parse_args()

    # Verify jarvis-web is running
    result = subprocess.run(
        ["systemctl", "--user", "is-active", "jarvis-web.service"],
        capture_output=True, text=True)
    if result.stdout.strip() != "active":
        print("ERROR: jarvis-web.service is not active.")
        print("Start it with: systemctl --user restart jarvis-web.service")
        sys.exit(1)

    try:
        import aiohttp  # noqa: F401
    except ImportError:
        print("ERROR: aiohttp not installed. Run: pip install aiohttp")
        sys.exit(1)

    try:
        summary_path = asyncio.run(run_test(args.output_dir))
        print(f"\nDone. Read the full report:\n  cat {summary_path}")
    except KeyboardInterrupt:
        print("\nInterrupted — cleaning up sentinel...")
        deactivate_sentinel()
        sys.exit(130)
    finally:
        deactivate_sentinel()


if __name__ == "__main__":
    main()
